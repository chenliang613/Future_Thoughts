#!/usr/bin/env python3
"""Generate AI Labor Market Impact Analysis PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xF4, 0x8C, 0x06)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
CARD_BG = RGBColor(0x24, 0x24, 0x3E)
SUBTLE_GRAY = RGBColor(0x99, 0x99, 0x99)


def add_dark_bg(slide):
    """Add dark background to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False, space_before=Pt(6), space_after=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_bottom_bar(slide, text="Anthropic Research | AI Labor Market Impact Analysis"):
    bar = add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), RGBColor(0x12, 0x12, 0x22))
    tb = add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.4))
    set_text(tb.text_frame, text, font_size=10, color=SUBTLE_GRAY, alignment=PP_ALIGN.LEFT)


# ========== Slide 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Decorative top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

tb = add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2))
set_text(tb.text_frame, "AI对劳动力市场的影响", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8))
set_text(tb.text_frame, "基于Anthropic \"观察暴露度\" 研究的深度分析报告", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Divider line
add_shape(slide, Inches(5), Inches(4.3), Inches(3.333), Inches(0.04), ACCENT_BLUE)

tb = add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.5))
set_text(tb.text_frame, "数据来源: Anthropic Research  |  O*NET  |  美国劳工统计局(BLS)", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5))
set_text(tb.text_frame, "2026年4月", font_size=16, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ========== Slide 2: Executive Summary ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "执行摘要", font_size=36, color=WHITE, bold=True)

# Four key finding cards
cards = [
    ("有限的当前影响", "AI实际应用远未达到理论能力\n实际覆盖率仅为可行应用的一小部分", ACCENT_BLUE),
    ("职业脆弱性", "暴露度较高的职业\n预计到2034年增长较慢", ACCENT_ORANGE),
    ("劳动者画像", "高暴露群体: 年龄偏大、女性居多\n受教育程度更高、薪资更高", ACCENT_GREEN),
    ("就业效应", "自2022年底以来未出现系统性失业\n但年轻劳动者招聘放缓明显", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.5 + i * 3.15)
    y = Inches(1.8)
    card = add_shape(slide, x, y, Inches(2.9), Inches(3.5), CARD_BG, color)
    # Color top accent
    accent = add_shape(slide, x, y, Inches(2.9), Inches(0.06), color)

    tb = add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.6))
    set_text(tb.text_frame, title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.5), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for line in desc.split('\n'):
        if tf.paragraphs[0].text == '':
            set_text(tf, line, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        else:
            add_paragraph(tf, line, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "核心发现: AI对劳动市场的影响正处于\"缓慢渗透\"阶段——理论能力远超实际部署，但年轻劳动者的招聘放缓是值得高度关注的早期信号。",
         font_size=15, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ========== Slide 3: Methodology ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "研究方法论: \"观察暴露度\" 指标", font_size=36, color=WHITE, bold=True)

# Subtitle
tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5))
set_text(tb.text_frame, "创新性地将LLM理论能力与真实使用数据相结合，衡量AI对职业的实际替代风险", font_size=16, color=LIGHT_GRAY)

# Three data sources
sources = [
    ("O*NET 职业数据库", "覆盖约800个美国职业\n详细编目各职业的任务构成", ACCENT_BLUE, "1"),
    ("Anthropic经济指数", "追踪Claude在专业场景中的\n实际使用模式和频率", ACCENT_ORANGE, "2"),
    ("Eloundou等(2023)理论评估", "使用β值(0-1)衡量\nLLM的理论能力边界", ACCENT_GREEN, "3"),
]

for i, (title, desc, color, num) in enumerate(sources):
    x = Inches(0.5 + i * 4.2)
    y = Inches(2.2)
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.2), CARD_BG, color)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.5), y - Inches(0.25), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = ctf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE
    run.font.bold = True

    tb = add_text_box(slide, x + Inches(0.2), y + Inches(0.4), Inches(3.4), Inches(0.5))
    set_text(tb.text_frame, title, font_size=17, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.4), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            set_text(tf, line, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        else:
            add_paragraph(tf, line, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Exposure scoring criteria
tb = add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.5))
set_text(tb.text_frame, "暴露度评分标准", font_size=22, color=WHITE, bold=True)

criteria_card = add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3), CARD_BG)

criteria = [
    ("理论可行性", "任务是否在LLM理论能力范围内"),
    ("实际使用验证", "是否有显著的工作相关使用记录"),
    ("自动化vs增强", "属于自动替代还是辅助增强模式"),
    ("任务覆盖占比", "在整体工作职责中的覆盖比例"),
]

for i, (title, desc) in enumerate(criteria):
    x = Inches(0.8 + i * 3.1)
    tb = add_text_box(slide, x, Inches(5.4), Inches(2.8), Inches(0.4))
    set_text(tb.text_frame, f"  {title}", font_size=15, color=ACCENT_BLUE, bold=True)
    tb = add_text_box(slide, x, Inches(5.8), Inches(2.8), Inches(0.4))
    set_text(tb.text_frame, f"  {desc}", font_size=12, color=LIGHT_GRAY)

add_bottom_bar(slide)

# ========== Slide 4: Coverage Gap ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "理论能力 vs 实际部署: 巨大落差", font_size=36, color=WHITE, bold=True)

tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5))
set_text(tb.text_frame, "各职业类别的AI覆盖率对比 — 实际应用远远落后于理论潜力", font_size=16, color=LIGHT_GRAY)

# Bar chart simulation
categories = [
    ("计算机与数学", 94, 33, ACCENT_BLUE),
    ("办公与行政", 78, 18, ACCENT_ORANGE),
    ("商业与金融", 72, 15, ACCENT_GREEN),
    ("法律", 68, 12, ACCENT_PURPLE),
    ("教育与培训", 55, 8, ACCENT_RED),
    ("体力与服务", 15, 2, SUBTLE_GRAY),
]

y_start = Inches(2.0)
bar_height = Inches(0.35)
gap = Inches(0.85)
label_width = Inches(2.5)
bar_area_width = Inches(8.5)

for i, (cat, theory, actual, color) in enumerate(categories):
    y = y_start + i * gap

    # Category label
    tb = add_text_box(slide, Inches(0.5), y - Inches(0.05), label_width, Inches(0.45))
    set_text(tb.text_frame, cat, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

    # Theory bar (dimmed)
    theory_w = Inches(8.5 * theory / 100)
    bar_bg = add_shape(slide, Inches(3.2), y, theory_w, bar_height, RGBColor(0x3A, 0x3A, 0x55))

    # Actual bar
    actual_w = Inches(8.5 * actual / 100)
    if actual > 0:
        bar_fg = add_shape(slide, Inches(3.2), y, actual_w, bar_height, color)

    # Percentage labels
    tb = add_text_box(slide, Inches(3.2) + theory_w + Inches(0.1), y - Inches(0.05), Inches(1.2), Inches(0.45))
    set_text(tb.text_frame, f"理论 {theory}%", font_size=11, color=SUBTLE_GRAY)

    if actual > 0:
        tb = add_text_box(slide, Inches(3.2) + actual_w + Inches(0.1), y + Inches(0.15), Inches(1.2), Inches(0.35))
        set_text(tb.text_frame, f"实际 {actual}%", font_size=11, color=color, bold=True)

# Key insight box
insight_card = add_shape(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7), CARD_BG, ACCENT_BLUE)
tb = add_text_box(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.6))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "关键洞察: 计算机与数学领域的实际覆盖率(33%)仅为理论潜力(94%)的1/3 — 约30%的劳动者所从事的体力和服务类工作几乎不受AI影响",
         font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ========== Slide 5: Most Exposed Occupations ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "最高暴露度职业分析", font_size=36, color=WHITE, bold=True)

# Top exposed occupations table
occupations = [
    ("计算机程序员", "75%", "自动化+增强", "高", ACCENT_RED),
    ("数据录入员", "67%", "主要自动化", "极高", ACCENT_RED),
    ("客服代表", "62%", "主要自动化", "极高", ACCENT_RED),
    ("技术文档撰写", "58%", "自动化+增强", "高", ACCENT_ORANGE),
    ("翻译与口译", "52%", "自动化+增强", "高", ACCENT_ORANGE),
    ("财务分析师", "45%", "主要增强", "中高", ACCENT_ORANGE),
    ("法律助理", "40%", "主要增强", "中高", ACCENT_ORANGE),
]

# Table header
header_y = Inches(1.6)
add_shape(slide, Inches(0.5), header_y, Inches(12.3), Inches(0.5), ACCENT_BLUE)
headers = [("职业名称", 0.7, 3.0), ("覆盖率", 4.0, 1.5), ("模式", 5.8, 2.2), ("替代风险", 8.3, 1.8)]
for text, x, w in headers:
    tb = add_text_box(slide, Inches(x), header_y + Inches(0.05), Inches(w), Inches(0.4))
    set_text(tb.text_frame, text, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (name, coverage, mode, risk, color) in enumerate(occupations):
    y = Inches(2.2 + i * 0.6)
    bg_color = CARD_BG if i % 2 == 0 else RGBColor(0x1E, 0x1E, 0x35)
    add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(0.55), bg_color)

    # Risk indicator
    add_shape(slide, Inches(0.5), y, Inches(0.06), Inches(0.55), color)

    row_data = [(name, 0.7, 3.0), (coverage, 4.0, 1.5), (mode, 5.8, 2.2), (risk, 8.3, 1.8)]
    for text, x, w in row_data:
        tb = add_text_box(slide, Inches(x), y + Inches(0.07), Inches(w), Inches(0.4))
        c = color if text == risk else LIGHT_GRAY
        b = True if text == coverage else False
        set_text(tb.text_frame, text, font_size=14, color=c, bold=b, alignment=PP_ALIGN.CENTER)

# Insight
tb = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5))
set_text(tf := tb.text_frame, "自动化主导型职业面临的替代风险最高，而增强型职业的AI更多扮演\"生产力放大器\"角色",
         font_size=15, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ========== Slide 6: Worker Demographics ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "高暴露度劳动者画像", font_size=36, color=WHITE, bold=True)

tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5))
set_text(tb.text_frame, "高暴露度群体 vs 低暴露度群体的关键差异 (ChatGPT发布前基线数据)", font_size=16, color=LIGHT_GRAY)

# Demographic comparison cards
demo_data = [
    ("女性比例", "+16%", "高暴露群体中女性占比\n高出16个百分点", ACCENT_PURPLE),
    ("平均收入", "+47%", "高暴露群体的平均收入\n高出47%", ACCENT_GREEN),
    ("研究生学历", "17.4% vs 4.5%", "高暴露群体的研究生\n比例是低暴露群体的近4倍", ACCENT_BLUE),
    ("亚裔占比", "近2倍", "高暴露群体中亚裔占比\n接近低暴露群体的两倍", ACCENT_ORANGE),
]

for i, (title, value, desc, color) in enumerate(demo_data):
    x = Inches(0.5 + i * 3.15)
    y = Inches(2.0)
    card = add_shape(slide, x, y, Inches(2.9), Inches(3.8), CARD_BG, color)
    add_shape(slide, x, y, Inches(2.9), Inches(0.06), color)

    tb = add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(2.6), Inches(0.5))
    set_text(tb.text_frame, title, font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    tb = add_text_box(slide, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(0.8))
    set_text(tb.text_frame, value, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_text_box(slide, x + Inches(0.15), y + Inches(2.1), Inches(2.6), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            set_text(tf, line, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        else:
            add_paragraph(tf, line, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Insight
insight_card = add_shape(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7), CARD_BG, ACCENT_GREEN)
tb = add_text_box(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.6))
set_text(tb.text_frame, "关键启示: AI首先影响的是高收入、高学历的知识工作者 — 这颠覆了\"技术首先取代低技能工人\"的传统认知",
         font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ========== Slide 7: Employment Impact ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_RED)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "就业影响: 现状与早期信号", font_size=36, color=WHITE, bold=True)

# Left panel: No systematic unemployment
left_card = add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.9), Inches(4.8), CARD_BG, ACCENT_GREEN)
add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.9), Inches(0.06), ACCENT_GREEN)

tb = add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.3), Inches(0.5))
set_text(tb.text_frame, "整体失业率: 暂无显著冲击", font_size=22, color=ACCENT_GREEN, bold=True)

unemployment_points = [
    "自ChatGPT发布以来，高暴露度和低暴露度\n群体之间未检测到系统性失业率差异",
    "研究框架可识别超过1个百分点的就业影响",
    "整体劳动市场保持相对稳定",
]
tb = add_text_box(slide, Inches(0.8), Inches(2.7), Inches(5.3), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
for j, point in enumerate(unemployment_points):
    if j == 0:
        set_text(tf, f"  {point}", font_size=14, color=LIGHT_GRAY)
    else:
        add_paragraph(tf, f"  {point}", font_size=14, color=LIGHT_GRAY, space_before=Pt(14))

# Right panel: Young worker warning
right_card = add_shape(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(4.8), CARD_BG, ACCENT_RED)
add_shape(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(0.06), ACCENT_RED)

tb = add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(0.5))
set_text(tb.text_frame, "年轻劳动者: 招聘放缓警报", font_size=22, color=ACCENT_RED, bold=True)

young_points = [
    "高暴露职业的求职成功率\n在ChatGPT后下降约14%",
    "低暴露职业维持稳定的2%月招聘率",
    "影响集中在22-25岁群体\n年长劳动者未观察到下降",
]
tb = add_text_box(slide, Inches(7.2), Inches(2.7), Inches(5.3), Inches(3.5))
tf = tb.text_frame
tf.word_wrap = True
for j, point in enumerate(young_points):
    if j == 0:
        set_text(tf, f"  {point}", font_size=14, color=LIGHT_GRAY)
    else:
        add_paragraph(tf, f"  {point}", font_size=14, color=LIGHT_GRAY, space_before=Pt(14))

# Key metric highlight
add_shape(slide, Inches(7.2), Inches(4.8), Inches(5.0), Inches(1.0), RGBColor(0x3A, 0x1A, 0x1A), ACCENT_RED)
tb = add_text_box(slide, Inches(7.4), Inches(4.9), Inches(4.6), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "14%", font_size=36, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "高暴露职业求职成功率下降幅度", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# ========== Slide 8: BLS Projections ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "BLS就业预测与暴露度的相关性", font_size=36, color=WHITE, bold=True)

tb = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5))
set_text(tb.text_frame, "美国劳工统计局(BLS)对2034年的就业增长预测验证了\"观察暴露度\"指标的有效性", font_size=16, color=LIGHT_GRAY)

# Core finding card
core = add_shape(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(2.5), CARD_BG, ACCENT_PURPLE)
add_shape(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(0.06), ACCENT_PURPLE)

tb = add_text_box(slide, Inches(2.5), Inches(2.5), Inches(8.3), Inches(0.5))
set_text(tb.text_frame, "核心发现", font_size=24, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(2.5), Inches(3.2), Inches(8.3), Inches(1.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "每增加10个百分点的观察暴露度", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "对应减少0.6个百分点的就业增长预测", font_size=28, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER, space_before=Pt(12))

# What this means
tb = add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.5))
set_text(tb.text_frame, "这意味着什么?", font_size=20, color=WHITE, bold=True)

meanings = [
    "BLS的独立预测与Anthropic的暴露度指标高度一致",
    "高暴露度职业预计未来10年将经历更慢的就业增长",
    "\"观察暴露度\"是衡量AI劳动市场影响的有效先行指标",
]
tb = add_text_box(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
for j, m in enumerate(meanings):
    if j == 0:
        set_text(tf, f"  {m}", font_size=15, color=LIGHT_GRAY)
    else:
        add_paragraph(tf, f"  {m}", font_size=15, color=LIGHT_GRAY, space_before=Pt(8))

add_bottom_bar(slide)

# ========== Slide 9: Limitations & Caveats ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_ORANGE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "研究局限与注意事项", font_size=36, color=WHITE, bold=True)

caveats = [
    ("方法论局限", "无法捕捉所有经济disruption渠道\n强调在显著效应出现前建立基准测量的重要性", ACCENT_ORANGE),
    ("数据采集偏差", "基于调查的工作转换数据可能低估实际就业变动\n初入劳动市场的年轻人缺少先前职业分类", ACCENT_RED),
    ("替代解释", "招聘放缓可能存在其他原因:\n继续深造、留在现有岗位、行业转换等", ACCENT_PURPLE),
    ("理论评估时效", "基于Eloundou等人2023年初的能力评估\n当前AI能力已显著超越当时水平", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(caveats):
    x = Inches(0.5 + (i % 2) * 6.4)
    y = Inches(1.6 + (i // 2) * 2.6)
    card = add_shape(slide, x, y, Inches(6.0), Inches(2.2), CARD_BG, color)
    add_shape(slide, x, y, Inches(0.08), Inches(2.2), color)

    tb = add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.4), Inches(0.5))
    set_text(tb.text_frame, title, font_size=20, color=color, bold=True)

    tb = add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        if j == 0:
            set_text(tf, line, font_size=14, color=LIGHT_GRAY)
        else:
            add_paragraph(tf, line, font_size=14, color=LIGHT_GRAY, space_before=Pt(6))

add_bottom_bar(slide)

# ========== Slide 10: Implications & Recommendations ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
set_text(tb.text_frame, "启示与建议", font_size=36, color=WHITE, bold=True)

# Three columns for different stakeholders
stakeholders = [
    ("对政策制定者", [
        "建立AI劳动市场影响的持续监测机制",
        "重点关注年轻劳动者就业入口问题",
        "在影响显现前制定预防性政策",
        "投资终身学习和技能转换项目",
    ], ACCENT_BLUE),
    ("对企业管理者", [
        "评估AI对核心岗位的暴露度水平",
        "制定\"增强优先\"的AI部署策略",
        "为高暴露岗位员工设计转型路径",
        "建立人机协作的最佳实践",
    ], ACCENT_GREEN),
    ("对劳动者个人", [
        "识别自身岗位的AI暴露风险等级",
        "发展AI难以替代的核心竞争力",
        "掌握AI工具以转为\"增强\"模式",
        "年轻求职者需关注行业暴露趋势",
    ], ACCENT_ORANGE),
]

for i, (title, items, color) in enumerate(stakeholders):
    x = Inches(0.5 + i * 4.2)
    y = Inches(1.5)
    card = add_shape(slide, x, y, Inches(3.8), Inches(4.8), CARD_BG, color)
    add_shape(slide, x, y, Inches(3.8), Inches(0.06), color)

    tb = add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.5))
    set_text(tb.text_frame, title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    tb = add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            set_text(tf, f"  {item}", font_size=13, color=LIGHT_GRAY)
        else:
            add_paragraph(tf, f"  {item}", font_size=13, color=LIGHT_GRAY, space_before=Pt(14))

add_bottom_bar(slide)

# ========== Slide 11: Conclusion ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

tb = add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0))
set_text(tb.text_frame, "总结", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Key takeaways
takeaways = [
    ("AI对劳动市场的影响处于早期阶段", "理论能力远超实际部署，但变化正在加速"),
    ("知识工作者首当其冲", "高收入、高学历群体面临最高暴露度，颠覆传统认知"),
    ("年轻劳动者是最敏感的风向标", "22-25岁群体招聘率下降14%，值得高度关注"),
    ("\"增强\"而非\"替代\"是主流模式", "AI更多扮演生产力工具角色，但自动化比例在上升"),
    ("现在是建立基准和采取行动的关键窗口期", "在影响全面显现前做好准备，方能化危为机"),
]

for i, (main, sub) in enumerate(takeaways):
    y = Inches(2.2 + i * 0.95)
    card = add_shape(slide, Inches(1.5), y, Inches(10.3), Inches(0.8), CARD_BG)
    add_shape(slide, Inches(1.5), y, Inches(0.06), Inches(0.8), ACCENT_BLUE)

    tb = add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(9.8), Inches(0.4))
    set_text(tb.text_frame, main, font_size=16, color=WHITE, bold=True)

    tb = add_text_box(slide, Inches(1.8), y + Inches(0.4), Inches(9.8), Inches(0.35))
    set_text(tb.text_frame, sub, font_size=13, color=LIGHT_GRAY)

# Source
tb = add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4))
set_text(tb.text_frame, "来源: Anthropic Research — anthropic.com/research/labor-market-impacts",
         font_size=12, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)

# Save
output_path = "/Users/apple/Future_Thoughts/AI影响劳动市场分析报告.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
