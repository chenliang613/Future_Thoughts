# -*- coding: utf-8 -*-
"""生成单页PPT：微软 AI 治理分析"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 配色 ----------
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x5F, 0x5F, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MSBLUE= RGBColor(0x00, 0x78, 0xD4)   # 微软品牌蓝
BG    = RGBColor(0xF4, 0xF7, 0xFB)
C1    = RGBColor(0x00, 0x78, 0xD4)   # 合规 蓝
C2    = RGBColor(0x10, 0x7C, 0x41)   # 承诺 绿
C3    = RGBColor(0x5C, 0x2D, 0x91)   # 体系 紫
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xDD, 0xE4, 0xEC)
FONT  = "PingFang SC"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def set_ea(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def rect(x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def text(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        for (t, sz, col, bold) in ln:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            set_ea(r)
    return tb

# ---------- 背景 ----------
rect(0, 0, prs.slide_width, prs.slide_height, BG)
rect(0, 0, prs.slide_width, Inches(0.18), MSBLUE)

# ---------- 标题 ----------
text(Inches(0.55), Inches(0.32), Inches(9.8), Inches(0.9),
     [[("微软 AI 治理分析", 30, INK, True)],
      [("“合规即竞争力” —— 以欧盟 AI 法案对齐与版权赔付承诺，构建云上 AI 信任壁垒", 14, SUB, False)]],
     sp_after=4)

tag = rect(Inches(10.45), Inches(0.42), Inches(2.35), Inches(0.55), MSBLUE)
tf = tag.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "治理 → 信任 → 分发"; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
set_ea(r)

# ---------- 核心主张 ----------
text(Inches(0.55), Inches(1.42), Inches(12.2), Inches(0.5),
     [[("核心打法：", 14, MSBLUE, True),
       ("把监管合规与法律风险“内化”为微软自身责任 —— 替客户扫清采用障碍，将治理能力沉淀为 Azure 的差异化护城河", 14, INK, False)]])

# ---------- 三张卡片 ----------
cards = [
    (C1, "① 合规先行", "对齐欧盟 AI 法案",
     [("每份客户合同中承诺遵守 EU AI Act 等全部适用法规", True),
      ("法案分期生效：2025.2 禁止性条款 → 2025.8 GPAI 义务", False),
      ("高风险义务经 Digital Omnibus 延至 2027.12", False),
      ("设跨职能工作组：治理+工程+法务+政策", False),
      ("欧盟数据驻留承诺，回应主权监管压力", False)]),
    (C2, "② 责任承诺", "客户版权承诺 (CCC)",
     [("2023.9 推出、10.1 生效，扩展商业 Copilot 赔付", True),
      ("客户因 Copilot 输出被诉版权侵权，微软出庭抗辩", False),
      ("并承担由此产生的判决与和解赔偿金额", False),
      ("前提：须开启内置护栏与内容过滤器", False),
      ("覆盖 M365 / GitHub Copilot；不含免费与消费版", False)]),
    (C3, "③ 治理体系", "把“负责任AI”产品化",
     [("Responsible AI Standard + 年度透明度报告", True),
      ("Azure AI Content Safety、透明度说明、内容凭证", False),
      ("Purview 合规管理器模板、风险评估与文档框架", False),
      ("Trust Center 设 EU AI Act 合规专区", False),
      ("AI 客户承诺(2024.6)：知识共享·工具·文档支持", False)]),
]

cx, cw, gap, cy, ch = Inches(0.55), Inches(3.97), Inches(0.18), Inches(2.05), Inches(3.62)
for i, (color, name, role, bullets) in enumerate(cards):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    rect(x, cy, cw, ch, CARD, line=LINE, line_w=Pt(1))
    rect(x, cy, cw, Inches(0.82), color)
    text(Emu(int(x)+Inches(0.22)), Emu(int(cy)+Inches(0.10)), Emu(int(cw)-Inches(0.4)), Inches(0.7),
         [[(name, 16, WHITE, True)], [(role, 11, WHITE, False)]], sp_after=1)
    lines = []
    for (t, hl) in bullets:
        lines.append([("▪ ", 12, color, True), (t, 11, INK if hl else SUB, hl)])
    text(Emu(int(x)+Inches(0.22)), Emu(int(cy)+Inches(0.98)), Emu(int(cw)-Inches(0.42)), Inches(2.5),
         lines, sp_after=6)

# ---------- 底部战略逻辑 ----------
fy = Inches(5.92)
rect(Inches(0.55), fy, Inches(12.23), Inches(1.18), RGBColor(0x1B, 0x29, 0x3A))
insights = [
    ("治理即信任", "合规对齐+赔付承诺，化解企业法律与合规顾虑"),
    ("信任即分发", "信任转化为云上采纳，放大 Azure AI 触达"),
    ("合规即壁垒", "把监管复杂度转为能力，中小厂商难以复制"),
    ("价值上移", "降低落地门槛，放大模型与行业应用价值"),
]
iw = Inches(3.0)
for i, (h, d) in enumerate(insights):
    x = Emu(int(Inches(0.72)) + i * int(iw))
    text(x, Emu(int(fy)+Inches(0.18)), Emu(int(iw)-Inches(0.16)), Inches(1.0),
         [[(h, 13.5, RGBColor(0x5E,0xB0,0xFF), True)], [(d, 10.5, RGBColor(0xDD,0xDD,0xDD), False)]], sp_after=3)

text(Inches(0.55), Inches(7.20), Inches(12.2), Inches(0.3),
     [[("资料来源：Microsoft Trust Center / On the Issues 博客 / 2025 Responsible AI Transparency Report；EU AI Act 官方时间线（2023.9–2026.06）", 8, SUB, False)]])

out = "/Users/apple/Future_Thoughts/微软AI治理分析.pptx"
prs.save(out)
print("saved:", out)
