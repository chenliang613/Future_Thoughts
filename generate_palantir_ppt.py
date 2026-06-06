#!/usr/bin/env python3
"""Generate Palantir 深度分析报告 — 重点：本体(Ontology) 与 AIP。

数据来源（Palantir Q1 FY2026 业绩 / 官方文档）：
- 营收 $1.6B，YoY +85%（上市以来最快）；US 营收首次 >100% 增长
- US 商业 $595M (+133% YoY, +18% QoQ)；US 政府 $687M (+84% YoY)
- TTM 商业客户 1,007 (+31% YoY)；GAAP 净利 $871M；调整后经营利润率 60%
- FY2026 指引上调至 $7.65–7.662B (+71%)；US 商业 >$3.224B (+120%)；ARR $6.5B+
- 本体：组织的运营层/数字孪生，语义(对象/属性/链接)+动力(动作/函数/动态安全)
- AIP：把生成式AI接入运营；安全LLM接入；AI Mesh(Foundry+Apollo+AIP)；Bootcamp 5天0到用例，~75%转化
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------- Color scheme ----------------
DARK_BG     = RGBColor(0x10, 0x14, 0x24)
PANEL_BG    = RGBColor(0x18, 0x1E, 0x33)
CARD_BG     = RGBColor(0x1F, 0x27, 0x42)
CARD_BG2    = RGBColor(0x24, 0x2E, 0x4F)
ACCENT      = RGBColor(0x2E, 0x9B, 0xF0)   # Palantir 蓝
ACCENT_DK   = RGBColor(0x16, 0x6A, 0xB8)
CYAN        = RGBColor(0x1A, 0xC8, 0xC0)
ORANGE      = RGBColor(0xF4, 0x8C, 0x06)
GOLD        = RGBColor(0xFF, 0xC4, 0x3D)
GREEN       = RGBColor(0x35, 0xCB, 0x7A)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE      = RGBColor(0x9B, 0x6BE if False else 0x59, 0xC6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT       = RGBColor(0xCF, 0xD7, 0xE6)
GRAY        = RGBColor(0x97, 0xA2, 0xB6)
DARKLINE    = RGBColor(0x30, 0x3A, 0x5C)

FONT = "Microsoft YaHei"   # PowerPoint 会以系统 CJK 字体回退 (Mac: PingFang)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = Inches(13.333), Inches(7.5)


# ---------------- Helpers ----------------
def slide_blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return s


def rect(s, x, y, w, h, fill, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of dict(text,size,color,bold,align,space_before,space_after)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if "space_before" in ln:
            p.space_before = Pt(ln["space_before"])
        if "space_after" in ln:
            p.space_after = Pt(ln["space_after"])
        if "line_spacing" in ln:
            p.line_spacing = ln["line_spacing"]
        runs = ln["text"] if isinstance(ln["text"], list) else [ln]
        for rr in runs:
            run = p.add_run()
            run.text = rr["text"]
            run.font.size = Pt(rr.get("size", ln.get("size", 14)))
            run.font.color.rgb = rr.get("color", ln.get("color", LIGHT))
            run.font.bold = rr.get("bold", ln.get("bold", False))
            run.font.name = FONT
    return tb


def header(s, kicker, title, accent=ACCENT):
    rect(s, 0, 0, SW, Inches(0.10), accent)
    rect(s, Inches(0.55), Inches(0.42), Inches(0.10), Inches(0.62), accent)
    txt(s, Inches(0.78), Inches(0.34), Inches(11.8), Inches(0.34),
        [{"text": kicker, "size": 12, "color": accent, "bold": True}])
    txt(s, Inches(0.78), Inches(0.60), Inches(11.8), Inches(0.55),
        [{"text": title, "size": 25, "color": WHITE, "bold": True}])


def page_no(s, n):
    txt(s, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
        [{"text": f"{n:02d}", "size": 11, "color": GRAY, "align": PP_ALIGN.RIGHT}])


def bullet_card(s, x, y, w, h, accent, title, items, title_size=15, body_size=11.5):
    rect(s, x, y, w, h, CARD_BG)
    rect(s, x, y, Inches(0.07), h, accent)
    lines = [{"text": title, "size": title_size, "color": WHITE, "bold": True,
              "space_after": 6}]
    for it in items:
        if isinstance(it, tuple):
            lines.append({"text": [{"text": "• ", "color": accent, "bold": True, "size": body_size},
                                   {"text": it[0], "color": WHITE, "bold": True, "size": body_size},
                                   {"text": it[1], "color": LIGHT, "size": body_size}],
                          "space_before": 5, "line_spacing": 1.05})
        else:
            lines.append({"text": [{"text": "• ", "color": accent, "bold": True, "size": body_size},
                                   {"text": it, "color": LIGHT, "size": body_size}],
                          "space_before": 5, "line_spacing": 1.05})
    txt(s, x + Inches(0.28), y + Inches(0.20), w - Inches(0.48), h - Inches(0.36), lines)


# =====================================================================
# Slide 1 — 封面
# =====================================================================
s = slide_blank()
rect(s, 0, 0, SW, SH, DARK_BG)
# 装饰条
rect(s, 0, 0, SW, Inches(0.14), ACCENT)
rect(s, 0, Inches(7.36), SW, Inches(0.14), ACCENT_DK)
# 左侧大色块
rect(s, Inches(0.0), Inches(2.0), Inches(0.18), Inches(3.4), CYAN)

txt(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5),
    [{"text": "企业 AI 应用层价值标杆 · 深度分析", "size": 16, "color": CYAN, "bold": True}])
txt(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(1.6),
    [{"text": [{"text": "Palantir ", "size": 52, "color": WHITE, "bold": True},
               {"text": "深度分析报告", "size": 44, "color": WHITE, "bold": True}]},
     {"text": "本体 (Ontology) × AIP — 解码企业智能化的价值飞轮",
      "size": 22, "color": ACCENT, "bold": True, "space_before": 10}])

# 三个要点条
yb = Inches(4.5)
for i, (t, c) in enumerate([
        ("本体 Ontology：组织的数字孪生与运营操作系统", CYAN),
        ("AIP：把生成式 AI 接入真实业务流程的引擎", ORANGE),
        ("Q1 FY2026 营收 +85%，应用层价值兑现的实证", GOLD)]):
    rect(s, Inches(0.9), yb + Inches(0.62) * i, Inches(0.14), Inches(0.40), c)
    txt(s, Inches(1.18), yb + Inches(0.62) * i, Inches(10.5), Inches(0.42),
        [{"text": t, "size": 15, "color": LIGHT, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)

txt(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4),
    [{"text": "行业智能化项目与产业发展分析  ·  2026-06  ·  数据截至 Palantir Q1 FY2026",
      "size": 12, "color": GRAY}])

# =====================================================================
# Slide 2 — 核心观点 / 执行摘要
# =====================================================================
s = slide_blank()
header(s, "EXECUTIVE SUMMARY", "核心观点：为什么 Palantir 是“倒金字塔”应用层的标杆", CYAN)

txt(s, Inches(0.78), Inches(1.40), Inches(11.8), Inches(0.6),
    [{"text": [{"text": "一句话结论：", "size": 15, "color": GOLD, "bold": True},
               {"text": "Palantir 不卖模型、不卖算力，而是用「本体」把企业数据资产变成可被 AI 操作的运营系统，"
                        "用 AIP 把大模型接入真实业务闭环——价值不在模型本身，而在“模型驱动业务行动”。",
                "size": 15, "color": LIGHT}], "line_spacing": 1.15}])

cards = [
    (CYAN, "本体是护城河", [
        "把“数据/模型”映射为业务“名词+动词”",
        "语义层(对象·属性·链接)+动力层(动作·函数)",
        "组织级数字孪生，AI 与人共享同一真相"]),
    (ORANGE, "AIP 是变现引擎", [
        "生成式 AI 不止于“聊天”，而是触发行动",
        "安全 LLM 接入，数据不外泄给第三方",
        "AI Mesh：Foundry+Apollo+AIP 全栈交付"]),
    (GOLD, "Bootcamp 是打法", [
        "5 天从 0 到真实用例，跑客户真实数据",
        "~75% 转化率，销售周期从~1年压缩到数天",
        "以“先见效再付费”重塑企业软件 GTM"]),
    (GREEN, "财务是实证", [
        "Q1 FY2026 营收 +85%，上市以来最快",
        "US 商业 +133%，调整后经营利润率 60%",
        "FY26 指引上调至 $7.65B+ (+71%)"]),
]
cx = Inches(0.55)
cw = Inches(3.0)
gap = Inches(0.10)
for i, (c, t, items) in enumerate(cards):
    x = cx + (cw + gap) * i
    bullet_card(s, x, Inches(2.30), cw, Inches(3.05), c, t, items,
                title_size=15, body_size=11)

txt(s, Inches(0.78), Inches(5.65), Inches(11.8), Inches(1.4),
    [{"text": "对本项目的意义", "size": 14, "color": ACCENT, "bold": True, "space_after": 4},
     {"text": "Palantir 证明了「倒金字塔」的可行性：在硬件与模型之上，应用层凭借“本体+工作流+反馈学习”形成"
              "复利护城河，捕获了远超基础设施的经济价值。其核心方法论——“先建组织本体、再把AI接入行动闭环、"
              "以Bootcamp快速验证ROI”——正是行业智能化项目成功的关键路径参照。",
      "size": 12.5, "color": LIGHT, "line_spacing": 1.18}])
page_no(s, 2)

# =====================================================================
# Slide 3 — 公司与业务全景
# =====================================================================
s = slide_blank()
header(s, "COMPANY OVERVIEW", "公司全景：从政府情报到企业 AI 的操作系统")

# 左：发展脉络
rect(s, Inches(0.55), Inches(1.45), Inches(5.7), Inches(5.5), CARD_BG)
rect(s, Inches(0.55), Inches(1.45), Inches(5.7), Inches(0.07), ACCENT)
txt(s, Inches(0.85), Inches(1.62), Inches(5.2), Inches(0.4),
    [{"text": "发展脉络与业务结构", "size": 16, "color": WHITE, "bold": True}])
timeline = [
    ("2003", "成立，以反恐/情报分析起家", CYAN),
    ("Gotham", "政府与国防：情报、作战、决策", ACCENT),
    ("Foundry", "商业：企业级数据运营平台", ORANGE),
    ("Apollo", "持续交付，跨云/本地/边缘部署", PURPLE),
    ("AIP (2023)", "把 LLM 接入运营，进入企业AI爆发期", GOLD),
    ("2026", "US 营收首破 100% 增长，AI 需求驱动", GREEN),
]
ty = Inches(2.20)
for i, (k, v, c) in enumerate(timeline):
    yy = ty + Inches(0.76) * i
    rect(s, Inches(0.95), yy + Inches(0.06), Inches(0.14), Inches(0.45), c)
    txt(s, Inches(1.25), yy, Inches(1.55), Inches(0.55),
        [{"text": k, "size": 13, "color": c, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.85), yy, Inches(3.25), Inches(0.55),
        [{"text": v, "size": 11.5, "color": LIGHT}], anchor=MSO_ANCHOR.MIDDLE)

# 右：四大平台
txt(s, Inches(6.55), Inches(1.50), Inches(6.3), Inches(0.4),
    [{"text": "四大平台 + 一个内核", "size": 16, "color": WHITE, "bold": True}])
plats = [
    ("Gotham", "面向政府/国防的运营分析平台", ACCENT),
    ("Foundry", "面向企业的数据集成与运营平台", ORANGE),
    ("Apollo", "持续交付层：任意环境自动部署", PURPLE),
    ("AIP", "AI 平台：把生成式AI接入运营闭环", GOLD),
]
for i, (t, d, c) in enumerate(plats):
    col = i % 2
    row = i // 2
    x = Inches(6.55) + (Inches(3.18)) * col
    y = Inches(2.05) + Inches(1.30) * row
    rect(s, x, y, Inches(3.05), Inches(1.15), CARD_BG2)
    rect(s, x, y, Inches(3.05), Inches(0.06), c)
    txt(s, x + Inches(0.2), y + Inches(0.16), Inches(2.7), Inches(0.4),
        [{"text": t, "size": 15, "color": c, "bold": True}])
    txt(s, x + Inches(0.2), y + Inches(0.58), Inches(2.7), Inches(0.5),
        [{"text": d, "size": 10.5, "color": LIGHT, "line_spacing": 1.05}])

# 内核条
rect(s, Inches(6.55), Inches(4.80), Inches(6.28), Inches(0.95), DARKLINE)
rect(s, Inches(6.55), Inches(4.80), Inches(0.10), Inches(0.95), CYAN)
txt(s, Inches(6.85), Inches(4.90), Inches(5.9), Inches(0.8),
    [{"text": [{"text": "共同内核 = 本体 Ontology  ", "size": 14, "color": CYAN, "bold": True},
               {"text": "所有平台都构建在同一套组织数字孪生之上，", "size": 11.5, "color": LIGHT}]},
     {"text": "这是 Palantir 区别于一般 BI / SaaS / 模型厂商的根本。",
      "size": 11.5, "color": LIGHT, "space_before": 3}])

# 关键数据条
rect(s, Inches(6.55), Inches(5.95), Inches(6.28), Inches(1.0), CARD_BG)
metrics = [("$1.6B", "Q1营收", GREEN), ("+85%", "YoY", GOLD),
           ("1,007", "商业客户", ACCENT), ("60%", "调整后\n经营利润率", CYAN)]
for i, (v, k, c) in enumerate(metrics):
    x = Inches(6.55) + Inches(1.57) * i
    txt(s, x, Inches(6.05), Inches(1.55), Inches(0.5),
        [{"text": v, "size": 20, "color": c, "bold": True, "align": PP_ALIGN.CENTER}])
    txt(s, x, Inches(6.55), Inches(1.55), Inches(0.4),
        [{"text": k, "size": 10, "color": GRAY, "align": PP_ALIGN.CENTER, "line_spacing": 0.9}])
page_no(s, 3)

# =====================================================================
# Slide 4 — 本体是什么 (概念)
# =====================================================================
s = slide_blank()
header(s, "ONTOLOGY · 本体 (1/3)", "本体是什么：组织的“数字孪生”与运营操作系统", CYAN)

txt(s, Inches(0.78), Inches(1.42), Inches(11.9), Inches(0.85),
    [{"text": [{"text": "定义　", "size": 15, "color": CYAN, "bold": True},
               {"text": "本体是一个架在企业数字资产(数据集·虚拟表·模型)之上的运营层，"
                        "把它们与现实世界的对应物连接起来——从物理设备，到客户订单、金融交易等业务概念。"
                        "它是组织的数字孪生：既包含“语义”，也包含“动力”。",
                "size": 13.5, "color": LIGHT}], "line_spacing": 1.2}])

# 三层堆叠图：现实世界 / 本体 / 数字资产
layers = [
    ("现实世界  REAL WORLD", "物理设备 · 人员 · 客户订单 · 金融交易 · 业务流程", ORANGE, CARD_BG2),
    ("本体  ONTOLOGY（运营层 / 数字孪生）",
     "语义：对象 Objects · 属性 Properties · 链接 Links（名词）   ＋   动力：动作 Actions · 函数 Functions · 动态安全（动词）",
     CYAN, RGBColor(0x12, 0x33, 0x3A)),
    ("数字资产  DIGITAL ASSETS", "数据集 Datasets · 虚拟表 · 机器学习模型 · 数据管道", ACCENT, CARD_BG2),
]
ly = Inches(2.55)
for i, (t, d, c, bg) in enumerate(layers):
    h = Inches(1.05) if i != 1 else Inches(1.25)
    rect(s, Inches(1.6), ly, Inches(10.1), h, bg, line=c, line_w=1.5)
    txt(s, Inches(1.9), ly + Inches(0.13), Inches(9.5), Inches(0.4),
        [{"text": t, "size": 15 if i == 1 else 14, "color": c, "bold": True}])
    txt(s, Inches(1.9), ly + Inches(0.55), Inches(9.5), Inches(0.6),
        [{"text": d, "size": 11.5, "color": LIGHT, "line_spacing": 1.05}])
    ly = ly + h + Inches(0.18)

# 左右箭头标注
txt(s, Inches(0.55), Inches(2.75), Inches(1.0), Inches(1.4),
    [{"text": "映射\n↕", "size": 16, "color": GRAY, "align": PP_ALIGN.CENTER, "line_spacing": 1.0}],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(11.85), Inches(4.4), Inches(1.2), Inches(1.4),
    [{"text": "↕\n绑定", "size": 16, "color": GRAY, "align": PP_ALIGN.CENTER, "line_spacing": 1.0}],
    anchor=MSO_ANCHOR.MIDDLE)

txt(s, Inches(0.78), Inches(6.7), Inches(11.9), Inches(0.6),
    [{"text": [{"text": "关键洞察　", "size": 12.5, "color": GOLD, "bold": True},
               {"text": "本体让“人”和“AI Agent”读到同一份业务真相——这是企业 AI 可信、可落地、可治理的前提。",
                "size": 12.5, "color": LIGHT}]}])
page_no(s, 4)

# =====================================================================
# Slide 5 — 本体的构成 (语义 / 动力 / 引擎)
# =====================================================================
s = slide_blank()
header(s, "ONTOLOGY · 本体 (2/3)", "本体的两大要素与一个引擎", CYAN)

cols = [
    (CYAN, "① 语义元素 Semantic", "组织的“名词”——结构", [
        ("对象 Object：", "实体或事件(如设备、订单)"),
        ("属性 Property：", "对象的特征字段"),
        ("链接 Link：", "对象之间的关系网络"),
        ("结果：", "把孤立数据织成业务知识图谱")]),
    (ORANGE, "② 动力元素 Kinetic", "组织的“动词”——行动", [
        ("动作 Action：", "捕获操作员输入/编排决策"),
        ("函数 Function：", "承载任意复杂业务逻辑"),
        ("动态安全：", "按权限/条件控制写入"),
        ("行动日志：", "记录决策溯源，供学习复盘")]),
    (GOLD, "③ 本体引擎 Engine", "让本体可运行、可进化", [
        ("规模查询：", "对数十亿对象实时检索"),
        ("动作编排：", "编排数万级动作执行"),
        ("反馈学习：", "持续吸收反馈、迭代优化"),
        ("写入编排：", "数据漏斗保证索引实时一致")]),
]
cw = Inches(3.93)
cx = Inches(0.55)
gap = Inches(0.17)
for i, (c, t, sub, items) in enumerate(cols):
    x = cx + (cw + gap) * i
    rect(s, x, Inches(1.55), cw, Inches(4.55), CARD_BG)
    rect(s, x, Inches(1.55), cw, Inches(0.65), c)
    txt(s, x + Inches(0.25), Inches(1.62), cw - Inches(0.4), Inches(0.5),
        [{"text": t, "size": 16, "color": DARK_BG, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.25), Inches(2.35), cw - Inches(0.5), Inches(0.4),
        [{"text": sub, "size": 12, "color": c, "bold": True}])
    lines = []
    for k, v in items:
        lines.append({"text": [{"text": "▸ ", "color": c, "bold": True, "size": 12.5},
                               {"text": k, "color": WHITE, "bold": True, "size": 12.5},
                               {"text": v, "color": LIGHT, "size": 12.5}],
                      "space_before": 9, "line_spacing": 1.1})
    txt(s, x + Inches(0.25), Inches(2.85), cw - Inches(0.45), Inches(3.1), lines)

rect(s, Inches(0.55), Inches(6.30), Inches(12.25), Inches(0.78), DARKLINE)
txt(s, Inches(0.85), Inches(6.42), Inches(11.7), Inches(0.6),
    [{"text": [{"text": "语义 = 让 AI“看懂”业务；  动力 = 让 AI“改变”业务；  引擎 = 让这套体系在生产环境实时运转。  ",
                "size": 13, "color": CYAN, "bold": True},
               {"text": "三者合一，本体才从“数据模型”升级为“运营操作系统”。",
                "size": 13, "color": LIGHT}]}], anchor=MSO_ANCHOR.MIDDLE)
page_no(s, 5)

# =====================================================================
# Slide 6 — 本体为何是护城河
# =====================================================================
s = slide_blank()
header(s, "ONTOLOGY · 本体 (3/3)", "本体为何是 Palantir 最深的护城河", CYAN)

# 左：企业AI落地的痛点 vs 本体解法
txt(s, Inches(0.78), Inches(1.45), Inches(6.0), Inches(0.4),
    [{"text": "企业 AI 落地的核心痛点 → 本体解法", "size": 16, "color": WHITE, "bold": True}])
pains = [
    ("数据孤岛、口径不一", "本体统一语义，单一业务真相"),
    ("LLM 会“幻觉”、不可信", "AI 在结构化对象/动作上行动，可约束可溯源"),
    ("AI 只能“建议”不能“执行”", "动作层把AI输出直接写回业务系统"),
    ("难治理、权限失控", "动态安全 + 行动日志，细粒度合规"),
    ("一次性项目、不复利", "反馈学习让系统越用越聪明"),
]
yy = Inches(2.0)
for i, (p, sol) in enumerate(pains):
    rect(s, Inches(0.78), yy, Inches(5.75), Inches(0.84), CARD_BG)
    rect(s, Inches(0.78), yy, Inches(0.07), Inches(0.84), CYAN)
    txt(s, Inches(1.0), yy + Inches(0.10), Inches(5.4), Inches(0.32),
        [{"text": [{"text": "痛点  ", "size": 11, "color": RED, "bold": True},
                   {"text": p, "size": 12.5, "color": WHITE, "bold": True}]}])
    txt(s, Inches(1.0), yy + Inches(0.45), Inches(5.4), Inches(0.32),
        [{"text": [{"text": "→  ", "size": 11.5, "color": GREEN, "bold": True},
                   {"text": sol, "size": 11.5, "color": LIGHT}]}])
    yy = yy + Inches(0.96)

# 右：护城河复利飞轮
txt(s, Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.4),
    [{"text": "复利护城河：越用越深", "size": 16, "color": WHITE, "bold": True}])
flywheel = [
    ("更多业务接入本体", "对象/链接更丰富", ACCENT),
    ("AI 在本体上做更多动作", "覆盖更多工作流", ORANGE),
    ("行动日志沉淀决策数据", "形成专有数据资产", GOLD),
    ("反馈学习持续优化", "效果提升、迁移成本上升", GREEN),
]
fy = Inches(2.05)
for i, (t, d, c) in enumerate(flywheel):
    rect(s, Inches(7.0), fy, Inches(5.5), Inches(0.82), CARD_BG2)
    rect(s, Inches(7.0), fy, Inches(0.55), Inches(0.82), c)
    txt(s, Inches(7.0), fy, Inches(0.55), Inches(0.82),
        [{"text": f"{i+1}", "size": 20, "color": DARK_BG, "bold": True, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.75), fy + Inches(0.10), Inches(4.6), Inches(0.7),
        [{"text": t, "size": 13.5, "color": WHITE, "bold": True},
         {"text": d, "size": 11, "color": LIGHT, "space_before": 2}])
    if i < 3:
        txt(s, Inches(9.5), fy + Inches(0.80), Inches(0.6), Inches(0.20),
            [{"text": "↓", "size": 14, "color": c, "align": PP_ALIGN.CENTER}])
    fy = fy + Inches(1.0)

rect(s, Inches(7.0), Inches(6.15), Inches(5.5), Inches(0.92), DARKLINE)
txt(s, Inches(7.25), Inches(6.26), Inches(5.0), Inches(0.7),
    [{"text": [{"text": "结论：", "size": 13, "color": CYAN, "bold": True},
               {"text": "护城河不在算法，而在“被本体建模的组织知识 + 沉淀的决策数据”——竞品难以复制，客户难以迁移。",
                "size": 11.5, "color": LIGHT, "line_spacing": 1.1}]}], anchor=MSO_ANCHOR.MIDDLE)
page_no(s, 6)

# =====================================================================
# Slide 7 — AIP 是什么
# =====================================================================
s = slide_blank()
header(s, "AIP · AI 平台 (1/3)", "AIP 是什么：把生成式 AI 接入“运营”的引擎", ORANGE)

txt(s, Inches(0.78), Inches(1.42), Inches(11.9), Inches(0.8),
    [{"text": [{"text": "定位　", "size": 15, "color": ORANGE, "bold": True},
               {"text": "AIP (Artificial Intelligence Platform) 把大模型从“会聊天”升级为“会干活”——"
                        "让 LLM 在本体的对象与动作之上感知数据、调用工具、触发真实业务行动，并在企业级安全边界内运行。",
                "size": 13.5, "color": LIGHT}], "line_spacing": 1.2}])

# 对比：通用聊天机器人 vs AIP
rect(s, Inches(0.78), Inches(2.45), Inches(5.75), Inches(2.0), CARD_BG)
rect(s, Inches(0.78), Inches(2.45), Inches(5.75), Inches(0.5), DARKLINE)
txt(s, Inches(0.98), Inches(2.50), Inches(5.4), Inches(0.4),
    [{"text": "通用聊天机器人", "size": 14, "color": GRAY, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
for i, t in enumerate(["脱离企业数据，易幻觉", "只能给“建议”，无法执行",
                        "无权限/审计，难以治理", "项目难以规模化、不复利"]):
    txt(s, Inches(1.05), Inches(3.05) + Inches(0.33) * i, Inches(5.3), Inches(0.3),
        [{"text": [{"text": "✕ ", "size": 12, "color": RED, "bold": True},
                   {"text": t, "size": 12, "color": LIGHT}]}])

rect(s, Inches(6.78), Inches(2.45), Inches(6.05), Inches(2.0), CARD_BG)
rect(s, Inches(6.78), Inches(2.45), Inches(6.05), Inches(0.5), ORANGE)
txt(s, Inches(6.98), Inches(2.50), Inches(5.6), Inches(0.4),
    [{"text": "Palantir AIP", "size": 14, "color": DARK_BG, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
for i, t in enumerate(["扎根本体，回答有据可查", "可直接触发动作、写回系统",
                        "继承动态安全与行动日志", "工作流复用 + 反馈学习复利"]):
    txt(s, Inches(7.05), Inches(3.05) + Inches(0.33) * i, Inches(5.6), Inches(0.3),
        [{"text": [{"text": "✓ ", "size": 12, "color": GREEN, "bold": True},
                   {"text": t, "size": 12, "color": LIGHT}]}])

# 三大能力
caps = [("数据感知", "理解本体对象/属性/链接，回答基于真实业务", ACCENT),
        ("工具调用", "调度函数与动作，编排跨系统决策流程", CYAN),
        ("反馈学习", "捕获专家反馈，体验“每天都在变好”", GOLD)]
for i, (t, d, c) in enumerate(caps):
    x = Inches(0.78) + Inches(4.08) * i
    rect(s, x, Inches(4.75), Inches(3.85), Inches(1.45), CARD_BG2)
    rect(s, x, Inches(4.75), Inches(3.85), Inches(0.06), c)
    txt(s, x + Inches(0.22), Inches(4.92), Inches(3.5), Inches(0.4),
        [{"text": t, "size": 15, "color": c, "bold": True}])
    txt(s, x + Inches(0.22), Inches(5.38), Inches(3.5), Inches(0.7),
        [{"text": d, "size": 11.5, "color": LIGHT, "line_spacing": 1.1}])

txt(s, Inches(0.78), Inches(6.45), Inches(11.9), Inches(0.5),
    [{"text": [{"text": "一句话　", "size": 12.5, "color": ORANGE, "bold": True},
               {"text": "本体提供“业务语言与行动接口”，AIP 提供“会用这套语言行动的 AI”——二者缺一不可。",
                "size": 12.5, "color": LIGHT}]}])
page_no(s, 7)

# =====================================================================
# Slide 8 — AIP 架构
# =====================================================================
s = slide_blank()
header(s, "AIP · AI 平台 (2/3)", "AIP 架构：本体之上的安全 AI 编排层", ORANGE)

# 架构分层（从上到下）
arch = [
    ("应用与交付层  AI Mesh", "Foundry + Apollo + AIP 协同：从 LLM Web 应用、视觉-语言移动应用，到嵌入式边缘 AI 全谱交付",
     ORANGE, RGBColor(0x33, 0x26, 0x12)),
    ("AI 编排层  Orchestration", "Agent / 工作流编排 · 工具(动作·函数)调用 · 人在回路审批 · 专家反馈捕获与学习",
     CYAN, RGBColor(0x12, 0x30, 0x33)),
    ("安全 LLM 接入层  Secure LLM Access", "经 Palantir 托管基础设施安全接入商用与开源模型；传输数据不被第三方模型提供商留存",
     GOLD, RGBColor(0x33, 0x2C, 0x12)),
    ("本体层  Ontology", "对象·属性·链接(语义) + 动作·函数·动态安全(动力)——AI 行动的统一底座",
     ACCENT, RGBColor(0x12, 0x22, 0x3A)),
    ("数据与基础设施层", "数据集成/管道 · 跨云·本地·边缘部署(Apollo) · 细粒度权限与审计",
     PURPLE, CARD_BG2),
]
ay = Inches(1.55)
hh = Inches(1.0)
for i, (t, d, c, bg) in enumerate(arch):
    rect(s, Inches(1.3), ay, Inches(10.7), hh, bg, line=c, line_w=1.4)
    rect(s, Inches(1.3), ay, Inches(0.12), hh, c)
    txt(s, Inches(1.6), ay + Inches(0.13), Inches(10.2), Inches(0.4),
        [{"text": t, "size": 14.5, "color": c, "bold": True}])
    txt(s, Inches(1.6), ay + Inches(0.52), Inches(10.2), Inches(0.42),
        [{"text": d, "size": 11, "color": LIGHT, "line_spacing": 1.0}])
    ay = ay + hh + Inches(0.12)

txt(s, Inches(0.45), Inches(1.55), Inches(0.8), Inches(5.0),
    [{"text": "应\n用\n\n↓\n\n数\n据", "size": 12, "color": GRAY, "align": PP_ALIGN.CENTER, "line_spacing": 0.95}],
    anchor=MSO_ANCHOR.MIDDLE)
page_no(s, 8)

# =====================================================================
# Slide 9 — AIP Bootcamp 商业模式
# =====================================================================
s = slide_blank()
header(s, "AIP · AI 平台 (3/3)", "AIP Bootcamp：重塑企业软件的 GTM 打法", ORANGE)

txt(s, Inches(0.78), Inches(1.42), Inches(11.9), Inches(0.55),
    [{"text": [{"text": "核心理念　", "size": 14, "color": ORANGE, "bold": True},
               {"text": "“用客户的真实数据，在几天内交付真实工作流”——先让客户看到价值，再谈采购。",
                "size": 13.5, "color": LIGHT}]}])

# 左：5天流程
txt(s, Inches(0.78), Inches(2.15), Inches(6.0), Inches(0.4),
    [{"text": "5 天：从 0 到真实用例", "size": 15, "color": WHITE, "bold": True}])
steps = [
    ("接入", "连接客户真实数据，构建/复用本体对象", ACCENT),
    ("建模", "定义关键工作流与动作，明确决策点", CYAN),
    ("构建", "用 AIP 搭建 AI 驱动的运营应用", ORANGE),
    ("见效", "跑出可量化 ROI 的真实用例", GREEN),
]
sy = Inches(2.65)
for i, (t, d, c) in enumerate(steps):
    rect(s, Inches(0.78), sy, Inches(5.85), Inches(0.86), CARD_BG)
    rect(s, Inches(0.78), sy, Inches(0.9), Inches(0.86), c)
    txt(s, Inches(0.78), sy, Inches(0.9), Inches(0.86),
        [{"text": f"D{i+1}", "size": 17, "color": DARK_BG, "bold": True, "align": PP_ALIGN.CENTER}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(1.85), sy + Inches(0.12), Inches(4.7), Inches(0.65),
        [{"text": t, "size": 13.5, "color": c, "bold": True},
         {"text": d, "size": 11, "color": LIGHT, "space_before": 2}])
    sy = sy + Inches(0.98)

# 右：成效指标
txt(s, Inches(7.0), Inches(2.15), Inches(5.8), Inches(0.4),
    [{"text": "可量化的商业成效", "size": 15, "color": WHITE, "bold": True}])
kpis = [("~75%", "Bootcamp 转化率", GREEN),
        ("~1年 → 数天", "销售周期压缩", GOLD),
        ("+22%", "US商业新客环比增速(Q4)", ACCENT),
        ("1,007", "TTM 商业客户 (+31% YoY)", CYAN)]
for i, (v, k, c) in enumerate(kpis):
    col = i % 2
    row = i // 2
    x = Inches(7.0) + Inches(2.92) * col
    y = Inches(2.65) + Inches(1.45) * row
    rect(s, x, y, Inches(2.78), Inches(1.30), CARD_BG2)
    rect(s, x, y, Inches(2.78), Inches(0.06), c)
    txt(s, x, y + Inches(0.22), Inches(2.78), Inches(0.55),
        [{"text": v, "size": 24, "color": c, "bold": True, "align": PP_ALIGN.CENTER}])
    txt(s, x + Inches(0.1), y + Inches(0.82), Inches(2.58), Inches(0.4),
        [{"text": k, "size": 11, "color": LIGHT, "align": PP_ALIGN.CENTER, "line_spacing": 1.0}])

rect(s, Inches(7.0), Inches(5.65), Inches(5.83), Inches(1.0), DARKLINE)
txt(s, Inches(7.25), Inches(5.76), Inches(5.4), Inches(0.8),
    [{"text": [{"text": "GTM 创新：", "size": 12.5, "color": ORANGE, "bold": True},
               {"text": "把“先见效再付费”产品化，用高触达交付换取高转化与短周期——"
                        "本质是用本体+AIP的“快速验证能力”作为最强销售武器。",
                "size": 11.5, "color": LIGHT, "line_spacing": 1.12}]}], anchor=MSO_ANCHOR.MIDDLE)
page_no(s, 9)

# =====================================================================
# Slide 10 — 本体 × AIP 协同飞轮
# =====================================================================
s = slide_blank()
header(s, "SYNERGY", "本体 × AIP：1+1 远大于 2 的价值飞轮", CYAN)

# 中心图：两个圆（本体、AIP）+ 中间协同
rect(s, Inches(1.2), Inches(2.0), Inches(3.6), Inches(2.6), RGBColor(0x12, 0x30, 0x38),
     line=CYAN, line_w=2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.2), Inches(2.15), Inches(3.6), Inches(0.5),
    [{"text": "本体 Ontology", "size": 17, "color": CYAN, "bold": True, "align": PP_ALIGN.CENTER}])
txt(s, Inches(1.45), Inches(2.7), Inches(3.1), Inches(1.8),
    [{"text": "提供业务语言与行动接口", "size": 12, "color": LIGHT, "align": PP_ALIGN.CENTER},
     {"text": "· 统一语义真相", "size": 11.5, "color": LIGHT, "space_before": 6},
     {"text": "· 可执行的动作层", "size": 11.5, "color": LIGHT, "space_before": 3},
     {"text": "· 动态安全与溯源", "size": 11.5, "color": LIGHT, "space_before": 3}])

rect(s, Inches(8.5), Inches(2.0), Inches(3.6), Inches(2.6), RGBColor(0x33, 0x26, 0x12),
     line=ORANGE, line_w=2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.5), Inches(2.15), Inches(3.6), Inches(0.5),
    [{"text": "AIP", "size": 17, "color": ORANGE, "bold": True, "align": PP_ALIGN.CENTER}])
txt(s, Inches(8.75), Inches(2.7), Inches(3.1), Inches(1.8),
    [{"text": "提供会行动的 AI", "size": 12, "color": LIGHT, "align": PP_ALIGN.CENTER},
     {"text": "· 数据感知问答", "size": 11.5, "color": LIGHT, "space_before": 6},
     {"text": "· 工具/动作编排", "size": 11.5, "color": LIGHT, "space_before": 3},
     {"text": "· 专家反馈学习", "size": 11.5, "color": LIGHT, "space_before": 3}])

# 中间协同
rect(s, Inches(5.0), Inches(2.35), Inches(3.3), Inches(1.9), GOLD,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(5.0), Inches(2.55), Inches(3.3), Inches(1.5),
    [{"text": "协同效应", "size": 16, "color": DARK_BG, "bold": True, "align": PP_ALIGN.CENTER},
     {"text": "AI 在可信业务底座上", "size": 12, "color": DARK_BG, "bold": True,
      "align": PP_ALIGN.CENTER, "space_before": 8},
     {"text": "安全地“理解→决策→执行→学习”", "size": 12, "color": DARK_BG, "bold": True,
      "align": PP_ALIGN.CENTER, "space_before": 4}])
txt(s, Inches(4.78), Inches(3.1), Inches(0.3), Inches(0.4),
    [{"text": "→", "size": 22, "color": CYAN, "bold": True, "align": PP_ALIGN.CENTER}])
txt(s, Inches(8.2), Inches(3.1), Inches(0.3), Inches(0.4),
    [{"text": "←", "size": 22, "color": ORANGE, "bold": True, "align": PP_ALIGN.CENTER}])

# 底部：结果产出
outs = [("可信", "回答与行动皆有据可溯源", GREEN),
        ("可执行", "AI 输出直接写回业务系统", ACCENT),
        ("可治理", "继承权限/审计/合规边界", PURPLE),
        ("会复利", "反馈学习让系统越用越强", GOLD)]
for i, (t, d, c) in enumerate(outs):
    x = Inches(0.78) + Inches(3.05) * i
    rect(s, x, Inches(5.05), Inches(2.85), Inches(1.5), CARD_BG)
    rect(s, x, Inches(5.05), Inches(2.85), Inches(0.06), c)
    txt(s, x, Inches(5.25), Inches(2.85), Inches(0.45),
        [{"text": t, "size": 17, "color": c, "bold": True, "align": PP_ALIGN.CENTER}])
    txt(s, x + Inches(0.18), Inches(5.78), Inches(2.5), Inches(0.65),
        [{"text": d, "size": 11.5, "color": LIGHT, "align": PP_ALIGN.CENTER, "line_spacing": 1.1}])

txt(s, Inches(0.78), Inches(6.75), Inches(11.9), Inches(0.45),
    [{"text": "这正是 Palantir 难以被“模型厂商 + BI 工具”组合替代的根本原因：护城河来自二者的耦合，而非单点能力。",
      "size": 12, "color": LIGHT, "align": PP_ALIGN.CENTER}])
page_no(s, 10)

# =====================================================================
# Slide 11 — 财务与增长实证
# =====================================================================
s = slide_blank()
header(s, "FINANCIALS · Q1 FY2026", "财务实证：应用层价值正在加速兑现", GREEN)

# 顶部大指标
big = [("$1.6B", "总营收", "YoY +85% · 上市以来最快", GREEN),
       ("+133%", "US 商业营收增速", "$595M · QoQ +18%", ORANGE),
       ("60%", "调整后经营利润率", "经营利润 $984M", CYAN),
       ("$871M", "GAAP 净利润", "调整后 EPS $0.33", GOLD)]
for i, (v, k, d, c) in enumerate(big):
    x = Inches(0.55) + Inches(3.13) * i
    rect(s, x, Inches(1.5), Inches(2.98), Inches(1.85), CARD_BG)
    rect(s, x, Inches(1.5), Inches(2.98), Inches(0.07), c)
    txt(s, x, Inches(1.72), Inches(2.98), Inches(0.6),
        [{"text": v, "size": 30, "color": c, "bold": True, "align": PP_ALIGN.CENTER}])
    txt(s, x, Inches(2.42), Inches(2.98), Inches(0.4),
        [{"text": k, "size": 13, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}])
    txt(s, x + Inches(0.1), Inches(2.85), Inches(2.78), Inches(0.4),
        [{"text": d, "size": 10.5, "color": GRAY, "align": PP_ALIGN.CENTER}])

# 中部：结构解读
rect(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(2.55), CARD_BG)
rect(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(0.07), ACCENT)
txt(s, Inches(0.8), Inches(3.75), Inches(5.5), Inches(0.4),
    [{"text": "增长结构：商业与政府双轮驱动", "size": 15, "color": WHITE, "bold": True}])
for i, (k, v, c) in enumerate([
        ("US 商业营收", "$595M  ·  +133% YoY  ·  +18% QoQ", ORANGE),
        ("US 政府营收", "$687M  ·  +84% YoY", ACCENT),
        ("US 总营收", "首次突破 100% 增长", GOLD),
        ("TTM 商业客户", "1,007 家  ·  +31% YoY", CYAN),
        ("ARR", "$6.5B+ ，AIP 驱动加速", GREEN)]):
    yy = Inches(4.25) + Inches(0.37) * i
    rect(s, Inches(0.8), yy + Inches(0.05), Inches(0.12), Inches(0.26), c)
    txt(s, Inches(1.05), yy, Inches(2.0), Inches(0.35),
        [{"text": k, "size": 12, "color": WHITE, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.05), yy, Inches(3.4), Inches(0.35),
        [{"text": v, "size": 12, "color": LIGHT}], anchor=MSO_ANCHOR.MIDDLE)

# 右：指引
rect(s, Inches(6.83), Inches(3.6), Inches(5.98), Inches(2.55), CARD_BG)
rect(s, Inches(6.83), Inches(3.6), Inches(5.98), Inches(0.07), GOLD)
txt(s, Inches(7.08), Inches(3.75), Inches(5.5), Inches(0.4),
    [{"text": "FY2026 全年指引（上调）", "size": 15, "color": WHITE, "bold": True}])
txt(s, Inches(7.08), Inches(4.3), Inches(5.5), Inches(1.8),
    [{"text": [{"text": "全年营收　", "size": 13, "color": GOLD, "bold": True},
               {"text": "$7.650B – $7.662B（+71% YoY）", "size": 13, "color": LIGHT}],
      "space_before": 2},
     {"text": [{"text": "US 商业　", "size": 13, "color": GOLD, "bold": True},
               {"text": "> $3.224B（增速 ≥ 120%）", "size": 13, "color": LIGHT}],
      "space_before": 10},
     {"text": [{"text": "信号　", "size": 13, "color": GOLD, "bold": True},
               {"text": "较前次指引中值上调约 10 个百分点，AI 需求是主驱动。", "size": 13,
                "color": LIGHT}], "space_before": 10, "line_spacing": 1.15}])

txt(s, Inches(0.55), Inches(6.35), Inches(12.3), Inches(0.7),
    [{"text": [{"text": "解读　", "size": 12.5, "color": GREEN, "bold": True},
               {"text": "高增长 + 高利润率同时实现，说明本体+AIP 不是“烧钱换增长”，而是高毛利、可复利的应用层价值——"
                        "这正是“倒金字塔”中应用层应有的经济回报。",
                "size": 12.5, "color": LIGHT, "line_spacing": 1.15}]}])
page_no(s, 11)

# =====================================================================
# Slide 12 — 对行业智能化的启示
# =====================================================================
s = slide_blank()
header(s, "IMPLICATIONS", "对“行业智能化”的方法论启示与关键路径", GOLD)

txt(s, Inches(0.78), Inches(1.40), Inches(11.9), Inches(0.5),
    [{"text": [{"text": "Palantir 给行业智能化项目的可复制路径：", "size": 14, "color": GOLD, "bold": True},
               {"text": "先建本体，再接 AI，以快速验证驱动规模化。", "size": 14, "color": LIGHT}]}])

paths = [
    ("01", "先建“行业本体”", "把行业的核心对象、关系、动作、规则建模为统一语义底座，让 AI 与人共享真相。"
                          "这是任何行业智能化的“地基工程”。", CYAN),
    ("02", "把 AI 接入“行动闭环”", "不要止步于问答/报表，要让 AI 直接触发审批、调度、下单等真实动作，"
                              "在闭环中产生可量化 ROI。", ORANGE),
    ("03", "安全与治理前置", "动态权限、审计、决策溯源必须内建，否则行业(金融/医疗/制造)无法规模化采用。", PURPLE),
    ("04", "以“快速验证”破局", "用类 Bootcamp 的方式，几天内跑出真实用例，先证明价值再规模采购，压缩决策周期。", GREEN),
    ("05", "构建反馈复利", "沉淀专有决策数据、引入反馈学习，让系统越用越强，形成应用层独有护城河。", GOLD),
    ("06", "捕获应用层价值", "价值不在模型/算力，而在“模型驱动业务行动”。应用层应主动定义并捕获这部分增量价值。", ACCENT),
]
for i, (n, t, d, c) in enumerate(paths):
    col = i % 3
    row = i // 3
    x = Inches(0.55) + Inches(4.15) * col
    y = Inches(2.10) + Inches(2.30) * row
    rect(s, x, y, Inches(3.95), Inches(2.10), CARD_BG)
    rect(s, x, y, Inches(3.95), Inches(0.06), c)
    txt(s, x + Inches(0.25), y + Inches(0.18), Inches(1.2), Inches(0.6),
        [{"text": n, "size": 26, "color": c, "bold": True}])
    txt(s, x + Inches(1.25), y + Inches(0.22), Inches(2.55), Inches(0.65),
        [{"text": t, "size": 14, "color": WHITE, "bold": True, "line_spacing": 1.0}])
    txt(s, x + Inches(0.25), y + Inches(0.95), Inches(3.5), Inches(1.05),
        [{"text": d, "size": 11, "color": LIGHT, "line_spacing": 1.12}])
page_no(s, 12)

# =====================================================================
# Slide 13 — 风险与挑战
# =====================================================================
s = slide_blank()
header(s, "RISKS & CAVEATS", "风险与挑战：标杆不等于可直接照搬", RED)

risks = [
    ("估值与预期", "高速增长已被高估值充分定价，任何减速都可能引发剧烈回调；增长可持续性需持续验证。", RED),
    ("高触达交付", "Bootcamp/前向部署工程师模式人力密集，规模化与毛利之间存在张力，复制到长尾客户有难度。", ORANGE),
    ("实施门槛", "本体建模需要深度业务理解与组织变革，落地周期与change management成本不可低估。", GOLD),
    ("竞争与替代", "云厂商(微软/AWS/Google)与模型厂商向应用层延伸，“本体”概念正被更多平台模仿。", PURPLE),
    ("客户集中与合规", "政府/大客户占比高，受预算周期、地缘政治与数据合规监管影响显著。", CYAN),
    ("对标可迁移性", "Palantir 优势源于长期积累的本体工程能力，后来者难以短期复制其工程与方法论壁垒。", ACCENT),
]
for i, (t, d, c) in enumerate(risks):
    col = i % 2
    row = i // 2
    x = Inches(0.55) + Inches(6.25) * col
    y = Inches(1.55) + Inches(1.78) * row
    rect(s, x, y, Inches(6.0), Inches(1.60), CARD_BG)
    rect(s, x, y, Inches(0.08), Inches(1.60), c)
    txt(s, x + Inches(0.3), y + Inches(0.18), Inches(5.5), Inches(0.4),
        [{"text": [{"text": "⚠ ", "size": 14, "color": c, "bold": True},
                   {"text": t, "size": 15, "color": WHITE, "bold": True}]}])
    txt(s, x + Inches(0.3), y + Inches(0.66), Inches(5.5), Inches(0.85),
        [{"text": d, "size": 11.5, "color": LIGHT, "line_spacing": 1.15}])
page_no(s, 13)

# =====================================================================
# Slide 14 — 总结
# =====================================================================
s = slide_blank()
rect(s, 0, 0, SW, SH, DARK_BG)
rect(s, 0, 0, SW, Inches(0.14), ACCENT)
rect(s, 0, Inches(7.36), SW, Inches(0.14), CYAN)

txt(s, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.5),
    [{"text": "CONCLUSION · 总结", "size": 14, "color": CYAN, "bold": True}])
txt(s, Inches(0.9), Inches(1.25), Inches(11.6), Inches(0.7),
    [{"text": "本体是地基，AIP 是引擎，飞轮是护城河", "size": 30, "color": WHITE, "bold": True}])

takeaways = [
    ("本体 Ontology", "把企业数据资产升级为“可被 AI 操作的运营操作系统”——语义让 AI 看懂业务，"
                     "动力让 AI 改变业务，引擎让它实时运转。这是 Palantir 最深的护城河。", CYAN),
    ("AIP", "把生成式 AI 从“会聊天”变成“会干活”：在本体之上安全地理解、决策、执行、学习，"
            "并以 Bootcamp 把“快速验证 ROI”变成最强销售武器。", ORANGE),
    ("价值兑现", "Q1 FY2026 营收 +85%、US 商业 +133%、调整后经营利润率 60%——高增长与高利润并存，"
              "印证应用层可以捕获最大经济价值。", GREEN),
    ("对本项目", "行业智能化的关键路径 = 先建行业本体 → 把 AI 接入行动闭环 → 安全治理前置 → "
             "快速验证破局 → 反馈复利 → 主动捕获应用层价值。", GOLD),
]
yy = Inches(2.35)
for i, (t, d, c) in enumerate(takeaways):
    rect(s, Inches(0.9), yy, Inches(11.5), Inches(1.0), CARD_BG)
    rect(s, Inches(0.9), yy, Inches(0.1), Inches(1.0), c)
    txt(s, Inches(1.25), yy, Inches(2.2), Inches(1.0),
        [{"text": t, "size": 16, "color": c, "bold": True}], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(3.5), yy + Inches(0.14), Inches(8.6), Inches(0.75),
        [{"text": d, "size": 12.5, "color": LIGHT, "line_spacing": 1.18}], anchor=MSO_ANCHOR.MIDDLE)
    yy = yy + Inches(1.10)

txt(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
    [{"text": "“硬件之上模型创造 10 倍价值，模型之上应用创造 100 倍价值”——Palantir 是这一“倒金字塔”叙事最有力的实证。",
      "size": 13, "color": ACCENT, "bold": True}])

# 保存
out = "/Users/apple/Future_Thoughts/palantir分析报告.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
