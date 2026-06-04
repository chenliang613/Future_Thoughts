#!/usr/bin/env python3
"""生成《行业智能化发展范式》PPT（v2 - 五条第一性原理）

内容与《行业智能化发展第一性原理.md》保持一致：
① 价值守恒  ② 行业知识密度  ③ 标杆-规模化路径  ④ 场景可批量复制  ⑤ 边际成本递减
并嵌入刷新后的《行业智能化第一性原理.jpg》信息图。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── 配色 ──────────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
CARD_BG     = RGBColor(0x24, 0x24, 0x3E)
INNER_BG    = RGBColor(0x16, 0x16, 0x2A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
SUBTLE_GRAY = RGBColor(0x99, 0x99, 0x99)

ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
GOAL = RGBColor(0x2E, 0xCC, 0x71)
RED  = RGBColor(0xE7, 0x4C, 0x3C)

# 五条原理配色（与信息图一致）
P = [RGBColor(0x2E, 0x86, 0xAB), RGBColor(0xA2, 0x3B, 0x72),
     RGBColor(0xF1, 0x8F, 0x01), RGBColor(0x3B, 0x8C, 0x6E),
     RGBColor(0x6A, 0x4C, 0x93)]

FONT = "Microsoft YaHei"


# ─── 通用辅助 ─────────────────────────────────────────────────────────────
def add_dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape(slide, l, t, w, h, fill_color, border_color=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _style_run(run, size, color, bold):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def set_text(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, anchor=None):
    tf.clear()
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _style_run(p.add_run(), size, color, bold)
    p.runs[0].text = text
    return p


def add_para(tf, text, size=15, color=WHITE, bold=False, align=PP_ALIGN.LEFT, before=4, after=2):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    _style_run(p.add_run(), size, color, bold)
    p.runs[0].text = text
    return p


def textbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def title_bar(slide, title, accent=ACCENT_BLUE):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), accent)
    tb = textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.9))
    set_text(tb.text_frame, title, size=32, color=WHITE, bold=True)


def bottom_bar(slide, idx):
    add_shape(slide, Inches(0), Inches(7.02), Inches(13.333), Inches(0.48), RGBColor(0x12, 0x12, 0x22))
    tb = textbox(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.4))
    set_text(tb.text_frame, "行业智能化发展范式  ·  从第一性原理提炼  ·  2026",
             size=10, color=SUBTLE_GRAY)
    tb2 = textbox(slide, Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.4))
    set_text(tb2.text_frame, str(idx), size=10, color=SUBTLE_GRAY, align=PP_ALIGN.RIGHT)


PAGE = [0]
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s)
    return s

def finish(slide):
    PAGE[0] += 1
    bottom_bar(slide, PAGE[0])


def principle_card(slide, x_in, y_in, w_in, h_in, num, title, sub, color, body, body_size=12):
    """绘制一张原理卡片：编号圆 + 标题 + 副标题 + 内嵌正文块"""
    add_shape(slide, Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in), CARD_BG, color)
    add_shape(slide, Inches(x_in), Inches(y_in), Inches(w_in), Inches(0.06), color)
    # 编号圆
    add_shape(slide, Inches(x_in + 0.2), Inches(y_in + 0.2), Inches(0.55), Inches(0.55),
              color, shape=MSO_SHAPE.OVAL)
    tb = textbox(slide, Inches(x_in + 0.2), Inches(y_in + 0.23), Inches(0.55), Inches(0.5))
    set_text(tb.text_frame, num, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 标题
    tb = textbox(slide, Inches(x_in + 0.85), Inches(y_in + 0.2), Inches(w_in - 1.0), Inches(0.6))
    set_text(tb.text_frame, title, size=16, color=color, bold=True)
    # 副标题
    tb = textbox(slide, Inches(x_in + 0.2), Inches(y_in + 0.85), Inches(w_in - 0.4), Inches(0.5))
    set_text(tb.text_frame, sub, size=12.5, color=WHITE, bold=True)
    # 正文块
    bh = h_in - 1.6
    add_shape(slide, Inches(x_in + 0.2), Inches(y_in + 1.45), Inches(w_in - 0.4), Inches(bh), INNER_BG)
    tf = textbox(slide, Inches(x_in + 0.32), Inches(y_in + 1.5), Inches(w_in - 0.64), Inches(bh - 0.1)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, body, size=body_size, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), ACCENT_BLUE)
tb = textbox(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.4))
set_text(tb.text_frame, "行业智能化发展范式", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb = textbox(s, Inches(1), Inches(3.45), Inches(11.3), Inches(0.8))
set_text(tb.text_frame, "从第一性原理看 AI 在各行各业落地的共同路径",
         size=24, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_shape(s, Inches(5.17), Inches(4.5), Inches(3), Inches(0.04), ACCENT_BLUE)
tb = textbox(s, Inches(1), Inches(4.95), Inches(11.3), Inches(0.5))
set_text(tb.text_frame, "剥离行业外壳后，任何智能化项目都绕不过的五条底层规律",
         size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
tb = textbox(s, Inches(1), Inches(5.7), Inches(11.3), Inches(0.5))
set_text(tb.text_frame, "源自《行业智能化发展第一性原理》  ·  2026 年版",
         size=13, color=SUBTLE_GRAY, align=PP_ALIGN.CENTER)
finish(s)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — 为什么需要一套范式（倒金字塔）
# ════════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "为什么需要一套发展范式？")
tb = textbox(s, Inches(0.7), Inches(1.25), Inches(12), Inches(0.7))
set_text(tb.text_frame,
         "金融风控、电商导购、制造质检、医疗诊断 …… 各行业 Agent 千差万别，"
         "但剥离行业外壳，它们遵循同一套底层规律。范式的使命，是把「正金字塔」翻转为「倒金字塔」。",
         size=15, color=LIGHT_GRAY)

add_shape(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(3.85), CARD_BG, RED)
add_shape(s, Inches(0.7), Inches(2.25), Inches(5.7), Inches(0.07), RED)
tb = textbox(s, Inches(0.95), Inches(2.45), Inches(5.2), Inches(0.6))
set_text(tb.text_frame, "❌ 正金字塔（不可持续）", size=19, color=RED, bold=True)
tf = textbox(s, Inches(0.95), Inches(3.15), Inches(5.2), Inches(2.8)).text_frame
tf.word_wrap = True
set_text(tf, "硬件  ×100  ← 拿走绝大部分价值", size=16, color=WHITE, bold=True)
add_para(tf, "模型  ×10   ← 仅为硬件的 1/10", size=15, color=LIGHT_GRAY, before=10)
add_para(tf, "应用  ×1    ← 又小一个数量级", size=14, color=SUBTLE_GRAY, before=10)
add_para(tf, "结果：产业结构倒挂，市场质疑「AI 泡沫」。", size=14, color=RED, before=18)

add_shape(s, Inches(6.9), Inches(2.25), Inches(5.7), Inches(3.85), CARD_BG, GOAL)
add_shape(s, Inches(6.9), Inches(2.25), Inches(5.7), Inches(0.07), GOAL)
tb = textbox(s, Inches(7.15), Inches(2.45), Inches(5.2), Inches(0.6))
set_text(tb.text_frame, "✅ 倒金字塔（健康生态）", size=19, color=GOAL, bold=True)
tf = textbox(s, Inches(7.15), Inches(3.15), Inches(5.2), Inches(2.8)).text_frame
tf.word_wrap = True
set_text(tf, "硬件  ×1", size=15, color=SUBTLE_GRAY, bold=True)
add_para(tf, "模型  ×10   ← 在硬件之上放大 10 倍", size=15, color=LIGHT_GRAY, before=10)
add_para(tf, "应用  ×100  ← 行业应用创造 100 倍价值", size=16, color=WHITE, bold=True, before=10)
add_para(tf, "范式目标：让应用层稳定捕获 100 倍价值——", size=14, color=GOAL, before=18)
add_para(tf, "这正是五条第一性原理共同支撑的终局。", size=14, color=GOAL, before=2)
finish(s)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — 五条第一性原理总览（嵌入刷新后的信息图）
# ════════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "共同范式：五条第一性原理")
tb = textbox(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.5))
set_text(tb.text_frame, "从大量行业实践中反向提炼——彼此正交、缺一不可。",
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

img = "/Users/apple/Future_Thoughts/行业智能化第一性原理.jpg"
if os.path.exists(img):
    # 原图 18:10 = 1.8；按可用高度 5.2in 适配，避免纵向溢出
    h_in = 5.2
    w_in = h_in * 18 / 10
    left = Inches((13.333 - w_in) / 2)
    s.shapes.add_picture(img, left, Inches(1.7), width=Inches(w_in), height=Inches(h_in))
else:
    tb = textbox(s, Inches(2), Inches(3.2), Inches(9), Inches(1))
    set_text(tb.text_frame, "[ 信息图缺失：行业智能化第一性原理.jpg ]",
             size=18, color=SUBTLE_GRAY, align=PP_ALIGN.CENTER)
finish(s)


# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — 原理一·二·三 详解
# ════════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "原理详解（上）：价值 · 知识 · 路径")
cards = [
    ("①", "价值守恒原理", "没有正 ROI 就没有真智能化", P[0],
     "AI 落地的本质是一次价值的流转：\n\n被替代认知劳动成本 + 新创红利\n>\n推理 + 部署 + 组织变革成本\n\n任何技术革命的可持续性，都建立在"
     "「创造价值 > 消耗成本」之上。"),
    ("②", "行业知识密度原理", "模型是地板，SOP 是天花板（本体论）", P[1],
     "行业级智能体 =\n通用大模型能力 × 行业知识密度\n（本体论 / Skill / SOP / 语料）\n\n模型能力增长放缓，"
     "而行业知识密度几乎无上限——竞争力正从「模型更强」转向「知识沉淀更深」。"),
    ("③", "标杆-规模化路径原理", "先打灯塔，再批量复制", P[2],
     "PoC（验可行）\n→ 灯塔项目（验 ROI）\n→ 规模复制（验产业化）\n\nTo B 扩张都遵循「先有标杆、再有规模」。"
     "灯塔把后续客户销售周期从 12 个月压到 3 个月。"),
]
for i, (num, t, sub, c, body) in enumerate(cards):
    principle_card(s, 0.7 + i * 4.05, 1.45, 3.85, 5.1, num, t, sub, c, body, body_size=11.5)
finish(s)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — 原理四·五 详解
# ════════════════════════════════════════════════════════════════════════════
s = new_slide()
title_bar(s, "原理详解（下）：复制 · 边际成本", accent=P[4])
cards = [
    ("④", "场景可批量复制原理", "标准化即规模化", P[3],
     "一个场景能否被真正「产业化」，取决于它能否被抽象成标准模板、在行业中可重复落地。\n\n"
     "选点优先级 = 场景价值 × 复制系数。\n\n能复制 1000 次的 60 分场景，"
     "胜过只能定制 1 次的 95 分场景。"),
    ("⑤", "边际成本递减原理", "第 N 次部署接近零成本", P[4],
     "软件业的根本经济学就建立在边际成本递减之上——这是 SaaS 估值高于传统服务业的根因。\n\n"
     "当 1 次开发能服务 1000 个客户、第 1001 个客户部署成本可忽略时，"
     "应用层捕获 100 倍价值才有「数学基础」。"),
]
for i, (num, t, sub, c, body) in enumerate(cards):
    principle_card(s, 0.7 + i * 6.15, 1.45, 5.85, 5.1, num, t, sub, c, body, body_size=13)
finish(s)


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — 一句话总结 + 终局
# ════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_shape(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), GOAL)
tb = textbox(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
set_text(tb.text_frame, "一句话总结", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(s, Inches(1.2), Inches(1.8), Inches(10.93), Inches(2.7), CARD_BG, GOAL)
tf = textbox(s, Inches(1.6), Inches(2.05), Inches(10.1), Inches(2.3)).text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "以正 ROI 为底线，用行业知识（本体论 / SOP）塑造灵魂；",
         size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_para(tf, "先打灯塔标杆跑通单点，再以场景标准化打开复制空间，",
         size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER, before=8)
add_para(tf, "最后用边际成本递减释放规模红利。",
         size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER, before=8)
add_para(tf, "任何行业、任何场景的智能化项目，本质都是这五条原理的不同组合与权重排列。",
         size=14, color=GOAL, align=PP_ALIGN.CENTER, before=16)

add_shape(s, Inches(1.2), Inches(5.0), Inches(10.93), Inches(1.3), RGBColor(0x14, 0x32, 0x26), GOAL)
tf = textbox(s, Inches(1.4), Inches(5.15), Inches(10.5), Inches(1.0)).text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "倒金字塔产业终局", size=20, color=GOAL, bold=True, align=PP_ALIGN.CENTER)
add_para(tf, "硬件 ×1    →    模型 ×10    →    行业应用 ×100",
         size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, before=8)
finish(s)


# ─── 保存 ─────────────────────────────────────────────────────────────────
out = "/Users/apple/Future_Thoughts/行业智能化发展范式.pptx"
prs.save(out)
print(f"已生成：{out}（共 {len(prs.slides._sldIdLst)} 页）")
