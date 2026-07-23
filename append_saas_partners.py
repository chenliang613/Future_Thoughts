# -*- coding: utf-8 -*-
"""向 Anthropic云策略.pptx 追加三页：Databricks / Snowflake / Salesforce 合作详析"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 配色 ----------
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x5F, 0x5F, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLAY  = RGBColor(0xD9, 0x77, 0x57)
CLAYD = RGBColor(0xB5, 0x5A, 0x3D)
BG    = RGBColor(0xF7, 0xF4, 0xEF)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xE3, 0xDD, 0xD3)
DARK  = RGBColor(0x2B, 0x2B, 0x2B)
GREY  = RGBColor(0xDD, 0xDD, 0xDD)
PANEL = RGBColor(0xEF, 0xE7, 0xDE)
FONT  = "PingFang SC"

PATH = "/Users/apple/Future_Thoughts/Anthropic云策略.pptx"
prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def set_ea(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def rect(slide, x, y, w, h, fill, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        if line_sp:
            p.line_spacing = line_sp
        for (t, sz, col, bold) in ln:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            set_ea(r)
    return tb

def section(slide, x, y, w, h, accent, title):
    rect(slide, x, y, w, h, CARD, line=LINE)
    rect(slide, x, y, Inches(0.1), h, accent)
    text(slide, Emu(int(x)+Inches(0.26)), Emu(int(y)+Inches(0.12)), Emu(int(w)-Inches(0.4)), Inches(0.4),
         [[(title, 14, accent, True)]])

def bullets(slide, x, y, w, items, accent, sz=11.5):
    lines = []
    for it in items:
        hl = it.startswith("*")
        t = it[1:] if hl else it
        lines.append([("▪ ", sz, accent, True), (t, sz, INK if hl else SUB, hl)])
    text(slide, x, y, w, Inches(3), lines, sp_after=8, line_sp=1.08)

def add_page(partner_cn, partner_en, accent, subtitle, chip, snapshot, integ, value, cases, ins_title, ins_text, source):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, BG)
    rect(s, 0, 0, SW, Inches(0.18), CLAY)
    # 标题
    text(s, Inches(0.55), Inches(0.32), Inches(9.6), Inches(0.95),
         [[("Anthropic × " + partner_cn, 30, INK, True)],
          [(subtitle, 14, SUB, False)]], sp_after=4)
    # 右上 chip
    tag = rect(s, Inches(10.5), Inches(0.42), Inches(2.3), Inches(0.55), accent)
    tf = tag.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = chip; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    set_ea(r)

    # ---- 左：合作快照 ----
    lx, ly, lw, lh = Inches(0.55), Inches(1.45), Inches(3.45), Inches(4.4)
    rect(s, lx, ly, lw, lh, DARK)
    rect(s, lx, ly, lw, Inches(0.62), accent)
    text(s, Emu(int(lx)+Inches(0.22)), Emu(int(ly)+Inches(0.10)), Emu(int(lw)-Inches(0.4)), Inches(0.5),
         [[("合作快照 · " + partner_en, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    ry = Emu(int(ly) + Inches(0.82))
    for label, val in snapshot:
        text(s, Emu(int(lx)+Inches(0.24)), ry, Emu(int(lw)-Inches(0.46)), Inches(0.62),
             [[(label, 10, RGBColor(0xB7,0xB1,0xA8), False)],
              [(val, 13, WHITE, True)]], sp_after=1, line_sp=1.0)
        ry = Emu(int(ry) + Inches(0.595))

    # ---- 右：三个区块 ----
    rx, rw = Inches(4.2), Inches(8.58)
    # 合作内容（整合）
    section(s, rx, Inches(1.45), rw, Inches(2.05), accent, "① 合作内容与技术整合")
    bullets(s, Emu(int(rx)+Inches(0.28)), Inches(1.95), Emu(int(rw)-Inches(0.5)), integ, accent)
    # 战略价值
    bw = Inches(4.2)
    section(s, rx, Inches(3.62), bw, Inches(2.23), accent, "② 战略价值")
    bullets(s, Emu(int(rx)+Inches(0.28)), Inches(4.12), Emu(int(bw)-Inches(0.5)), value, accent, sz=11)
    # 客户案例
    cx2 = Emu(int(rx) + int(bw) + Inches(0.18))
    cw2 = Emu(int(rw) - int(bw) - Inches(0.18))
    section(s, cx2, Inches(3.62), cw2, Inches(2.23), accent, "③ 客户与案例")
    bullets(s, Emu(int(cx2)+Inches(0.28)), Inches(4.12), Emu(int(cw2)-Inches(0.5)), cases, accent, sz=11)

    # ---- 底部洞察条 ----
    iy = Inches(6.02)
    rect(s, Inches(0.55), iy, Inches(12.23), Inches(0.98), DARK)
    rect(s, Inches(0.55), iy, Inches(0.12), Inches(0.98), accent)
    text(s, Inches(0.88), Emu(int(iy)+Inches(0.12)), Inches(11.7), Inches(0.78),
         [[(ins_title + "：", 12.5, CLAY, True), (ins_text, 12, WHITE, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_sp=1.08)
    # 来源
    text(s, Inches(0.55), Inches(7.14), Inches(12.2), Inches(0.3),
         [[(source, 8, SUB, False)]])
    return s

# 品牌色
DBX = RGBColor(0xE8, 0x40, 0x2A)   # Databricks 红
SNO = RGBColor(0x29, 0xB5, 0xE8)   # Snowflake 蓝
SFC = RGBColor(0x00, 0xA1, 0xE0)   # Salesforce 蓝

# ============================================================
# 页：Databricks
# ============================================================
add_page(
    "Databricks", "Databricks", DBX,
    "五年期里程碑协议 —— 让 Claude 在企业「治理数据」上直接跑 Agent",
    "5 年 · 约 $100M",
    [("协议类型", "战略合作（里程碑级）"),
     ("公布时间", "2025.03.26"),
     ("期限 / 金额", "5 年 · 约 $100M"),
     ("触达客户", "10,000+ 企业"),
     ("覆盖云", "AWS / Azure / GCP"),
     ("荣誉", "2026 ISV「AI 转型年度伙伴」")],
    ["*Claude 原生上架 Data Intelligence Platform",
     "进入 Agent Bricks / AI Playground / Model Serving",
     "*经 Unity Catalog & Unity AI Gateway 统一治理调用",
     "Unity Catalog 的 SQL/Python 函数作为工具注入 Claude SDK",
     "前沿模型快速跟进上线（3.7 Sonnet → Fable 5）"],
    ["*在「自有专有数据」上构建领域专用 Agent，数据不出治理边界",
     "Constitutional AI 安全 × Unity 治理 = 受治理的企业级 Agent",
     "把 Agent 价值锚定在企业最敏感数据上 → 高粘性、高 Token 消耗",
     "与 Snowflake 形成「双数据平台」双覆盖"],
    ["*Block（原 Square）：以该合作为引擎，",
     "  驱动其开源 AI Agent「codename goose」",
     "面向 Databricks 上万家企业客户的存量数据资产",
     "数据科学 / 工程团队的 Agent 开发场景"],
    "战略本质",
    "数据在哪，智能就去哪。Databricks 出「治理 + 数据 + 客户」，Claude 出「智能」，Anthropic 借此把 Agent 牢牢嵌入企业专有数据之上。",
    "资料来源：Databricks 官方新闻稿与博客、TechTarget、Technology Magazine、NAND Research（2025.03–2026）；金额为公开报道口径。",
)

# ============================================================
# 页：Snowflake
# ============================================================
add_page(
    "Snowflake", "Snowflake", SNO,
    "$200M 多年期协议 + 联合 GTM —— 用「受治理的 AI」打入大型企业生产",
    "多年 · $200M",
    [("协议类型", "扩展战略合作"),
     ("期限 / 金额", "多年期 · $200M"),
     ("触达客户", "12,600+ 全球企业"),
     ("覆盖云", "Bedrock / Vertex / Azure"),
     ("联合 GTM", "是（联合销售团队）"),
     ("主题", "Governed AI 受治理 AI")],
    ["*Claude 驱动 Snowflake Cortex Code（编程）",
     "*Claude 驱动 Snowflake Intelligence（智能分析）",
     "上架 Claude Marketplace，面向安全的开发工作流",
     "跨 Bedrock / Vertex AI / Azure 三大云均可调用",
     "围绕「企业生产级工作负载」深度协同"],
    ["*主题即卖点：Governed AI —— 直击企业合规与可治理诉求",
     "联合 GTM：面向「全球最大企业」批量部署 AI Agent",
     "12,600+ 客户即时可达，分发杠杆极高",
     "与 Databricks 同时供货 = 通吃数据平台赛道算力"],
    ["*12,600+ Snowflake 企业客户为即时市场",
     "需求侧：企业对「合规、可治理」AI 的需求快速上升",
     "金融、医疗等数据敏感行业的生产化落地",
     "Cortex 生态内的 Agent / 代码 / 分析场景"],
    "战略本质",
    "在数据平台战争中两边都供货 —— 无论企业选 Snowflake 还是 Databricks，智能层都是 Claude。Anthropic 以「中立军火商」之姿吃下整条赛道。",
    "资料来源：Anthropic / Snowflake 官方新闻稿、TechTarget、BigDATAwire（2025–2026）；金额为公开披露口径。",
)

# ============================================================
# 页：Salesforce
# ============================================================
add_page(
    "Salesforce", "Salesforce", SFC,
    "以「信任边界」内置拿下受监管行业 —— Claude 成 Agentforce 首选模型",
    "受监管行业首选",
    [("协议类型", "扩展战略合作"),
     ("公布时间", "2025.10.14"),
     ("核心定位", "Agentforce 360 基础模型"),
     ("投资方", "Salesforce Ventures（C→G 轮）"),
     ("目标行业", "金融/医疗/安全/生命科学"),
     ("里程碑", "首个进入信任边界的 LLM")],
    ["*成为 Agentforce 受监管行业「首选模型」",
     "*全部 Claude 流量限定在 Salesforce VPC / 信任边界内",
     "Slack × MCP：Claude 直读 Slack/CRM 上下文并执行动作",
     "Salesforce 全球工程团队部署 Claude Code（反向成为客户）",
     "共建行业方案：先做金融（Claude for Financial Services）"],
    ["*「信任边界」内置 = 攻克受监管行业最大顾虑（安全/合规）",
     "搭车高增长引擎：Agentforce ARR $800M（FY26Q4，YoY+169%）",
     "双向绑定：既是供应商（进 Agentforce），又是客户（用 Claude Code）",
     "受监管行业是「100× 应用层价值」最厚、门槛最高之处"],
    ["*CrowdStrike：在 Agentforce 中用 Claude 构建 AI 体验",
     "*RBC Wealth Management：财富管理场景落地",
     "金融机构：理解金融工具、保险理赔与行业框架",
     "经 Slack 触达海量企业协作场景"],
    "战略本质",
    "受监管行业价值最厚但门槛是「信任」。Anthropic 用「进入对方信任边界」换取行业准入，把合规劣势转化为与 Salesforce 共建的壁垒。",
    "资料来源：Anthropic / Salesforce 官方新闻稿与投资者公告、Salesforce Ben、ETIH（2025.10–2026）。",
)

prs.save(PATH)
print("saved:", PATH, "| total slides:", len(prs.slides._sldIdLst))
