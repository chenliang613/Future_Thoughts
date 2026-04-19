#!/usr/bin/env python3
"""Generate single-page PPT: Agent时代数据基础设施 — 三层架构"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xF4, 0x8C, 0x06)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
ACCENT_CYAN = RGBColor(0x1A, 0xBC, 0x9C)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG = RGBColor(0x24, 0x24, 0x3E)
CARD_BG_2 = RGBColor(0x1E, 0x1E, 0x35)
SUBTLE_GRAY = RGBColor(0x99, 0x99, 0x99)
LAYER_INFRA = RGBColor(0x1C, 0x2B, 0x3A)    # 算力层底色
LAYER_BRAIN = RGBColor(0x1C, 0x2D, 0x38)    # 大脑层底色
LAYER_BIZ = RGBColor(0x1E, 0x25, 0x3A)      # 业务层底色
CONNECTOR_BG = RGBColor(0x2A, 0x1F, 0x3A)   # 连接层底色


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape(slide, left, top, width, height, fill_color, border_color=None,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_size=14, color=LIGHT_GRAY, bold=False,
                  space_before=Pt(4), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_small_card(slide, x, y, w, h, title, desc, accent_color):
    """Add a small card with top accent bar."""
    add_shape(slide, x, y, w, h, CARD_BG, accent_color)
    add_shape(slide, x, y, w, Inches(0.05), accent_color)
    # Title
    tb = add_text_box(slide, x + Inches(0.12), y + Inches(0.10), w - Inches(0.24), Inches(0.30))
    set_text(tb.text_frame, title, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    tb = add_text_box(slide, x + Inches(0.10), y + Inches(0.38), w - Inches(0.20), h - Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    set_text(tf, desc, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ========== Single Slide ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

# Title
tb = add_text_box(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.55))
set_text(tb.text_frame, "Agent时代的数据基础设施", font_size=30, color=WHITE,
         bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
tb = add_text_box(slide, Inches(0.5), Inches(0.60), Inches(12.3), Inches(0.40))
set_text(tb.text_frame,
         "数据是Agent的世界模型 -- Agent在没有人工干预的情况下，理解数据并做出正确决策的效率，决定了数据的真正价值",
         font_size=13, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# =====================================================================
#  Layout constants for three layers + connector band
# =====================================================================
margin_x = Inches(0.35)
content_w = Inches(12.6)

layer3_y = Inches(1.10)     # Top: 智能化业务层
layer3_h = Inches(1.55)

conn_y = Inches(2.72)       # Connector band
conn_h = Inches(0.72)

layer2_y = Inches(3.52)     # Middle: 基础设施大脑层
layer2_h = Inches(1.85)

layer1_y = Inches(5.45)     # Bottom: 算力层
layer1_h = Inches(1.15)

# =====================================================================
#  Layer 3 (Top): 智能化业务层 — Agent驱动
# =====================================================================
add_shape(slide, margin_x, layer3_y, content_w, layer3_h, LAYER_BIZ, ACCENT_PURPLE)

# Layer label (left side vertical-style label)
tb = add_text_box(slide, margin_x + Inches(0.10), layer3_y + Inches(0.08),
                  Inches(1.60), Inches(0.35))
set_text(tb.text_frame, "智能化业务层", font_size=14, color=ACCENT_PURPLE, bold=True)

tb = add_text_box(slide, margin_x + Inches(0.10), layer3_y + Inches(0.38),
                  Inches(1.60), Inches(0.25))
set_text(tb.text_frame, "Agent as New User", font_size=10, color=ACCENT_PURPLE, bold=False)

# Business scenario cards
biz_cards = [
    ("智能客服 Agent", "多轮对话 · 知识检索\n自动工单 · 情感分析", ACCENT_PURPLE),
    ("供应链 Agent", "需求预测 · 库存优化\n物流调度 · 风险预警", ACCENT_CYAN),
    ("金融风控 Agent", "实时反欺诈 · 信贷审批\n合规审查 · 市场监控", ACCENT_ORANGE),
    ("研发效能 Agent", "代码生成 · 测试自动化\nCI/CD编排 · 文档生成", ACCENT_GREEN),
    ("医疗诊断 Agent", "辅助诊断 · 影像分析\n药物推荐 · 病历摘要", ACCENT_RED),
    ("数据分析 Agent", "自然语言查询 · 自动建模\n报表生成 · 异常检测", ACCENT_BLUE),
]

biz_x0 = margin_x + Inches(1.80)
biz_y0 = layer3_y + Inches(0.10)
biz_card_w = Inches(1.72)
biz_card_h = Inches(1.35)
biz_gap = Inches(0.08)

for i, (title, desc, color) in enumerate(biz_cards):
    x = biz_x0 + i * (biz_card_w + biz_gap)
    add_small_card(slide, x, biz_y0, biz_card_w, biz_card_h, title, desc, color)

# =====================================================================
#  Connector Band: MCP / A2A / Skill / Harness
# =====================================================================
add_shape(slide, margin_x, conn_y, content_w, conn_h, CONNECTOR_BG, ACCENT_GOLD)

# Label
tb = add_text_box(slide, margin_x + Inches(0.15), conn_y + Inches(0.05),
                  content_w - Inches(0.30), Inches(0.22))
set_text(tb.text_frame,
         "Agent 连接协议层  —  打通「基础设施大脑」与「智能化业务」的桥梁",
         font_size=11, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Four connector items
connectors = [
    ("MCP", "Model Context Protocol", "标准化上下文接入\n让Agent发现并调用数据源与工具", ACCENT_BLUE),
    ("A2A", "Agent-to-Agent Protocol", "Agent间协作通信协议\n多Agent编排、任务分发与结果聚合", ACCENT_CYAN),
    ("Skill", "可复用能力单元", "封装领域知识为可调用技能\nAgent按需组合、即插即用", ACCENT_GREEN),
    ("Harness", "执行控制框架", "Agent运行时管控与编排\n权限·审计·限流·回滚保障", ACCENT_ORANGE),
]

conn_card_w = Inches(2.95)
conn_card_h = Inches(0.40)
conn_x0 = margin_x + Inches(0.15)
conn_gap = Inches(0.10)

for i, (name, full_name, desc, color) in enumerate(connectors):
    x = conn_x0 + i * (conn_card_w + conn_gap)
    y = conn_y + Inches(0.28)

    card = add_shape(slide, x, y, conn_card_w, conn_card_h, CARD_BG, color)
    # Name badge
    badge_w = Inches(0.55)
    add_shape(slide, x + Inches(0.05), y + Inches(0.06), badge_w, Inches(0.28), color)
    tb = add_text_box(slide, x + Inches(0.05), y + Inches(0.06), badge_w, Inches(0.28))
    set_text(tb.text_frame, name, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Full name + desc (two lines right of badge)
    tb = add_text_box(slide, x + Inches(0.65), y + Inches(0.03), conn_card_w - Inches(0.75), Inches(0.18))
    set_text(tb.text_frame, full_name, font_size=9, color=WHITE, bold=True)
    tb = add_text_box(slide, x + Inches(0.65), y + Inches(0.20), conn_card_w - Inches(0.75), Inches(0.20))
    set_text(tb.text_frame, desc.replace("\n", " | "), font_size=8, color=LIGHT_GRAY)

# Up/Down arrows between connector band and layers
for ax in [Inches(3.5), Inches(6.6), Inches(9.7)]:
    # Up arrow to biz layer
    arrow_up = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, ax, conn_y - Inches(0.12),
                                      Inches(0.25), Inches(0.15))
    arrow_up.fill.solid()
    arrow_up.fill.fore_color.rgb = ACCENT_GOLD
    arrow_up.line.fill.background()
    # Down arrow to brain layer
    arrow_dn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, ax, conn_y + conn_h - Inches(0.03),
                                      Inches(0.25), Inches(0.15))
    arrow_dn.fill.solid()
    arrow_dn.fill.fore_color.rgb = ACCENT_GOLD
    arrow_dn.line.fill.background()

# =====================================================================
#  Layer 2 (Middle): 基础设施大脑层 — 模型 + 数据
# =====================================================================
add_shape(slide, margin_x, layer2_y, content_w, layer2_h, LAYER_BRAIN, ACCENT_BLUE)

# Layer label
tb = add_text_box(slide, margin_x + Inches(0.10), layer2_y + Inches(0.08),
                  Inches(1.70), Inches(0.35))
set_text(tb.text_frame, "基础设施大脑层", font_size=14, color=ACCENT_BLUE, bold=True)

tb = add_text_box(slide, margin_x + Inches(0.10), layer2_y + Inches(0.38),
                  Inches(1.70), Inches(0.25))
set_text(tb.text_frame, "模型 + 数据", font_size=11, color=ACCENT_BLUE, bold=False)

# Two sub-sections side by side: 模型能力 | 数据能力
brain_inner_x = margin_x + Inches(1.80)
brain_inner_w = Inches(10.65)
brain_half_w = Inches(5.25)
brain_inner_y = layer2_y + Inches(0.10)
brain_inner_h = layer2_h - Inches(0.18)

# --- Left: 模型能力 ---
model_cards = [
    ("大语言模型", "GPT/Claude/Llama\n推理·生成·规划", ACCENT_PURPLE),
    ("多模态模型", "视觉·语音·代码\n跨模态理解", ACCENT_CYAN),
    ("专业模型", "领域微调·小模型\n低延迟·高精度", ACCENT_GREEN),
]

mc_x0 = brain_inner_x
mc_y0 = brain_inner_y + Inches(0.05)
mc_w = Inches(1.68)
mc_h = Inches(1.65)
mc_gap = Inches(0.07)

# Section title
tb = add_text_box(slide, mc_x0, mc_y0, brain_half_w, Inches(0.25))
set_text(tb.text_frame, "模型能力", font_size=11, color=ACCENT_CYAN, bold=True,
         alignment=PP_ALIGN.CENTER)

for i, (title, desc, color) in enumerate(model_cards):
    x = mc_x0 + i * (mc_w + mc_gap)
    add_small_card(slide, x, mc_y0 + Inches(0.28), mc_w, mc_h - Inches(0.35), title, desc, color)

# --- Right: 数据能力 ---
data_cards = [
    ("多模态存储", "向量库·图谱·时序库\n结构化+非结构化统一", ACCENT_BLUE),
    ("语义数据层", "元数据·数据目录\n业务语义·数据卡片", ACCENT_ORANGE),
    ("RAG与检索", "混合检索·Rerank\n实时召回·低延迟", ACCENT_GREEN),
    ("记忆系统", "长短期记忆·跨会话\n上下文沉淀与压缩", ACCENT_PURPLE),
    ("安全治理", "身份认证·行级权限\n审计追溯·反馈闭环", ACCENT_RED),
]

dc_x0 = brain_inner_x + brain_half_w + Inches(0.15)
dc_y0 = brain_inner_y + Inches(0.05)
dc_w = Inches(1.05)
dc_h = Inches(1.65)
dc_gap = Inches(0.05)

# Section title
tb = add_text_box(slide, dc_x0, dc_y0, brain_half_w, Inches(0.25))
set_text(tb.text_frame, "数据能力", font_size=11, color=ACCENT_ORANGE, bold=True,
         alignment=PP_ALIGN.CENTER)

for i, (title, desc, color) in enumerate(data_cards):
    x = dc_x0 + i * (dc_w + dc_gap)
    add_small_card(slide, x, dc_y0 + Inches(0.28), dc_w, dc_h - Inches(0.35), title, desc, color)

# =====================================================================
#  Layer 1 (Bottom): 算力层 — 基础设施硬件
# =====================================================================
add_shape(slide, margin_x, layer1_y, content_w, layer1_h, LAYER_INFRA, ACCENT_GREEN)

# Layer label
tb = add_text_box(slide, margin_x + Inches(0.10), layer1_y + Inches(0.08),
                  Inches(1.60), Inches(0.35))
set_text(tb.text_frame, "算力基础设施层", font_size=14, color=ACCENT_GREEN, bold=True)

tb = add_text_box(slide, margin_x + Inches(0.10), layer1_y + Inches(0.38),
                  Inches(1.60), Inches(0.25))
set_text(tb.text_frame, "Hardware Infra", font_size=10, color=ACCENT_GREEN, bold=False)

# Infra items
infra_items = [
    ("GPU/TPU集群", "训练·推理算力\nNVIDIA/AMD/自研芯片", ACCENT_GREEN),
    ("高速网络", "RDMA · InfiniBand\n低延迟互联", ACCENT_CYAN),
    ("分布式存储", "对象存储·块存储\n高吞吐IO", ACCENT_BLUE),
    ("云原生平台", "K8s · Serverless\n弹性伸缩·调度编排", ACCENT_PURPLE),
    ("边缘计算", "端侧推理·IoT网关\n低延迟本地决策", ACCENT_ORANGE),
    ("能源与散热", "绿色数据中心\n液冷·能效优化", ACCENT_RED),
]

inf_x0 = margin_x + Inches(1.80)
inf_y0 = layer1_y + Inches(0.10)
inf_w = Inches(1.72)
inf_h = Inches(0.95)
inf_gap = Inches(0.08)

for i, (title, desc, color) in enumerate(infra_items):
    x = inf_x0 + i * (inf_w + inf_gap)
    add_small_card(slide, x, inf_y0, inf_w, inf_h, title, desc, color)

# =====================================================================
#  Bottom Insight Bar
# =====================================================================
insight_y = Inches(6.70)
insight_card = add_shape(slide, margin_x, insight_y, content_w, Inches(0.70),
                         CARD_BG, ACCENT_BLUE)
tb = add_text_box(slide, margin_x + Inches(0.15), insight_y + Inches(0.05),
                  content_w - Inches(0.30), Inches(0.30))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf,
         "核心洞察：Agent取代人类成为数据的主要消费者 —— 数据就是Agent的世界模型",
         font_size=13, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, margin_x + Inches(0.15), insight_y + Inches(0.35),
                  content_w - Inches(0.30), Inches(0.30))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf,
         "数据的价值不再取决于你拥有多少，而取决于Agent在没有人工干预的情况下，理解数据并做出正确决策的效率 —— 谁让Agent低成本理解数据，谁就掌握智能化入口",
         font_size=10, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/apple/Future_Thoughts/Agent时代数据基础设施.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
