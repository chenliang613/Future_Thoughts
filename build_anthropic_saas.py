# -*- coding: utf-8 -*-
"""生成多页PPT：Anthropic 生态分析 —— 与 SaaS 厂商的合作策略"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 配色 ----------
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x5F, 0x5F, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLAY  = RGBColor(0xD9, 0x77, 0x57)   # Anthropic 品牌橙
CLAYD = RGBColor(0xB5, 0x5A, 0x3D)
BG    = RGBColor(0xF7, 0xF4, 0xEF)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xE3, 0xDD, 0xD3)
DARK  = RGBColor(0x2B, 0x2B, 0x2B)
BLUE  = RGBColor(0x42, 0x85, 0xF4)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
GOLD  = RGBColor(0xC8, 0x96, 0x2E)
PURP  = RGBColor(0x7C, 0x5C, 0xBF)
GREY  = RGBColor(0xDD, 0xDD, 0xDD)
FONT  = "PingFang SC"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def set_ea(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def rect(slide, x, y, w, h, fill, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2, line_sp=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        if line_sp:
            p.line_spacing = line_sp
        for (t, sz, col, bold) in ln:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            set_ea(r)
    return tb

def new_slide(title, subtitle, idx, total):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, BG)
    rect(s, 0, 0, SW, Inches(0.16), CLAY)
    text(s, Inches(0.55), Inches(0.34), Inches(11.2), Inches(0.9),
         [[(title, 26, INK, True)],
          [(subtitle, 13, SUB, False)]], sp_after=3)
    # 页码
    text(s, Inches(12.0), Inches(0.42), Inches(0.9), Inches(0.4),
         [[(f"{idx:02d} / {total:02d}", 11, CLAY, True)]], align=PP_ALIGN.RIGHT)
    rect(s, Inches(0.55), Inches(1.30), Inches(12.23), Pt(1.4), LINE)
    return s

def x_of(base, w, gap, i):
    return Emu(int(base) + i * (int(w) + int(gap)))

TOTAL = 8

# ============================================================
# 封面
# ============================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, 0, Inches(0.32), SH, CLAY)
rect(s, 0, Inches(4.55), SW, Pt(1.4), RGBColor(0x4A,0x4A,0x4A))
text(s, Inches(0.9), Inches(0.85), Inches(6), Inches(0.5),
     [[("行业智能化 · 产业生态研究", 14, CLAY, True)]])
text(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(2.2),
     [[("Anthropic 生态分析", 52, WHITE, True)],
      [("与 SaaS 厂商的合作策略", 34, RGBColor(0xE7,0xC8,0xB6), True)]], sp_after=8)
text(s, Inches(0.92), Inches(4.85), Inches(11.4), Inches(1.4),
     [[("以「模型即智能层」嵌入软件生态，以 MCP 定义连接标准，", 16, GREY, False)],
      [("用 Token 消耗捕获价值 —— 同时向应用层攀升的「合作 × 竞争」双重博弈", 16, GREY, False)]],
     sp_after=6)
text(s, Inches(0.92), Inches(6.7), Inches(11), Inches(0.4),
     [[("行业智能化项目组    |    2026-06    |    数据截至 2026 Q1", 11, RGBColor(0x9A,0x9A,0x9A), False)]])

# ============================================================
# 2. 战略全景 —— 模型即智能层 / 倒金字塔
# ============================================================
s = new_slide("战略命题：从「模型层」向「应用层」要价值", "在倒金字塔产业结构下，谁掌握应用层与连接标准，谁就掌握 AI 产业的价值主导权", 1, TOTAL)

# 左：倒金字塔
px, py, pw = Inches(0.55), Inches(1.65), Inches(5.3)
pyr = [
    (Inches(0.0),  pw,            CLAY,  "行业应用层 (SaaS)", "目标 100× 价值 · Anthropic 既供货又下场"),
    (Inches(1.15), Inches(3.6),   CLAYD, "模型层 (Claude)", "10× 价值 · Anthropic 主战场"),
    (Inches(2.05), Inches(1.9),   DARK,  "算力 / 硬件层", "1× · 见《Anthropic 云策略》"),
]
ty = py
for off, w, col, name, desc in pyr:
    bx = Emu(int(px) + int((pw - w) / 2))
    rect(s, bx, ty, w, Inches(0.95), col)
    text(s, bx, Emu(int(ty)+Inches(0.14)), w, Inches(0.7),
         [[(name, 14, WHITE, True)], [(desc, 9.5, RGBColor(0xF0,0xE8,0xE0), False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp_after=1)
    ty = Emu(int(ty) + Inches(1.08))
text(s, px, Emu(int(ty)+Inches(0.05)), pw, Inches(0.5),
     [[("健康生态 = 越往上价值越大；模型厂商必须借 SaaS 把模型变成「结果」", 10.5, SUB, False)]],
     align=PP_ALIGN.CENTER)

# 右：三句话命题
rx = Inches(6.3)
props = [
    ("①  模型是「成分」，不是终点", CLAY,
     "Claude 单卖 API 只是中间品；只有被 SaaS 封装成具体业务结果，才能触达 100× 的应用层价值。"),
    ("②  分发杠杆在 SaaS 手里", BLUE,
     "ServiceNow、Salesforce、Snowflake 等握有存量企业客户与数据；借其分发，远快于 Anthropic 自建销售。"),
    ("③  标准即护城河", GREEN,
     "通过开源 MCP 定义「模型↔软件」连接标准，把生态主导权握在手中，而非沦为可替换的算力供应商。"),
]
ry = Inches(1.65)
for h, col, d in props:
    rect(s, rx, ry, Inches(6.45), Inches(1.55), CARD, line=LINE)
    rect(s, rx, ry, Inches(0.12), Inches(1.55), col)
    text(s, Emu(int(rx)+Inches(0.32)), Emu(int(ry)+Inches(0.18)), Inches(6.0), Inches(1.2),
         [[(h, 15, INK, True)], [(d, 11.5, SUB, False)]], sp_after=4, line_sp=1.05)
    ry = Emu(int(ry) + Inches(1.72))

# ============================================================
# 3. 合作版图 —— 生态全景
# ============================================================
s = new_slide("合作版图：Claude 已成为企业软件的「智能成分」", "横跨数据平台、业务 SaaS、编程工具与横向应用，Claude 作为底层智能被广泛嵌入", 2, TOTAL)

cats = [
    (BLUE,  "数据 / AI 平台", "治理数据上的 Agent 底座",
     ["Snowflake Cortex（$200M）", "Databricks（5 年期）", "Palantir AIP", "IBM watsonx"]),
    (GREEN, "业务 SaaS（成分品牌）", "嵌入工作流的智能引擎",
     ["Salesforce / Agentforce", "ServiceNow", "Intuit · Docusign", "Notion · Canva · Asana"]),
    (CLAY,  "编程赛道（最大收入源）", "高 Token 消耗的杀手级场景",
     ["GitHub Copilot", "Cursor · Windsurf", "Replit · Vercel", "Claude Code（自营）"]),
    (PURP,  "横向 / 入口", "面向终端的新交互入口",
     ["Microsoft Copilot", "Amazon（Alexa+）", "Zoom · Slack", "Perplexity 等"]),
]
cx, cw, gap, cy, ch = Inches(0.55), Inches(2.95), Inches(0.18), Inches(1.62), Inches(3.05)
for i, (col, name, role, items) in enumerate(cats):
    x = x_of(cx, cw, gap, i)
    rect(s, x, cy, cw, ch, CARD, line=LINE)
    rect(s, x, cy, cw, Inches(0.82), col)
    text(s, Emu(int(x)+Inches(0.18)), Emu(int(cy)+Inches(0.10)), Emu(int(cw)-Inches(0.3)), Inches(0.7),
         [[(name, 13, WHITE, True)], [(role, 9.5, RGBColor(0xF2,0xF2,0xF2), False)]], sp_after=1)
    lines = [[("● ", 11, col, True), (it, 11.5, INK, False)] for it in items]
    text(s, Emu(int(x)+Inches(0.18)), Emu(int(cy)+Inches(0.98)), Emu(int(cw)-Inches(0.34)), Inches(1.9),
         lines, sp_after=9)

# 底部数据条
fy = Inches(4.95)
rect(s, Inches(0.55), fy, Inches(12.23), Inches(1.55), DARK)
stats = [
    ("$9B+", "年化收入（2025 末），约 80% 来自 API / 企业"),
    ("18%", "Claude Code 在编程市场份额，9 个月增 6 倍"),
    ("12,600+", "经 Snowflake 一家即可触达的企业客户"),
    (">50%", "Claude Code 收入来自企业（Netflix/Spotify/KPMG…）"),
]
iw = Inches(3.0)
for i, (big, d) in enumerate(stats):
    x = x_of(Inches(0.78), iw, 0, i)
    text(s, x, Emu(int(fy)+Inches(0.22)), Emu(int(iw)-Inches(0.2)), Inches(1.2),
         [[(big, 28, CLAY, True)], [(d, 10.5, GREY, False)]], sp_after=4, line_sp=1.0)

# ============================================================
# 4. 三种合作模式
# ============================================================
s = new_slide("三种合作模式：嵌入 · 联合分发 · 标准共建", "从「卖 Token」到「定标准」，合作层级越高，锁定越深、可替代性越低", 3, TOTAL)

models = [
    (CLAY, "模式一 · 成分品牌嵌入", "Model-as-Ingredient",
     "SaaS 把 Claude 作为底层智能封装进自有产品（Agentforce、Cortex），用户无感知。",
     "价值捕获：Token 消耗（用量计费）",
     "代表：Salesforce · ServiceNow · Notion"),
    (BLUE, "模式二 · 联合 GTM 分发", "Co-Sell & Marketplace",
     "与平台签多年期协议、设联合销售团队，借其存量客户与 Marketplace 放量。",
     "价值捕获：承诺采购额 + 渠道放量",
     "代表：Snowflake $200M · Databricks 5 年"),
    (GREEN, "模式三 · 标准共建", "MCP / Agentic Foundation",
     "开源 MCP 定义模型↔工具连接标准，捐入 Linux 基金会，掌握生态话语权。",
     "价值捕获：标准主导权 + 默认入口",
     "代表：MCP 10,000+ 服务器 · 全行业采纳"),
]
mx, mw, mgap, my, mh = Inches(0.55), Inches(3.97), Inches(0.18), Inches(1.65), Inches(3.45)
for i, (col, name, en, desc, val, rep) in enumerate(models):
    x = x_of(mx, mw, mgap, i)
    rect(s, x, my, mw, mh, CARD, line=LINE)
    rect(s, x, my, mw, Inches(0.92), col)
    text(s, Emu(int(x)+Inches(0.22)), Emu(int(my)+Inches(0.13)), Emu(int(mw)-Inches(0.4)), Inches(0.8),
         [[(name, 15, WHITE, True)], [(en, 10.5, RGBColor(0xF2,0xF2,0xF2), False)]], sp_after=1)
    text(s, Emu(int(x)+Inches(0.22)), Emu(int(my)+Inches(1.08)), Emu(int(mw)-Inches(0.44)), Inches(1.3),
         [[(desc, 12, INK, False)]], sp_after=2, line_sp=1.1)
    # 价值/代表
    text(s, Emu(int(x)+Inches(0.22)), Emu(int(my)+Inches(2.45)), Emu(int(mw)-Inches(0.44)), Inches(0.9),
         [[("◆ ", 10, col, True), (val, 11, INK, True)],
          [("▸ ", 10, col, True), (rep, 10.5, SUB, False)]], sp_after=4, line_sp=1.05)

# 升级箭头注释
text(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.4),
     [[("锁定深度 / 不可替代性：  成分嵌入  ──▶  联合分发  ──▶  标准共建（最深）", 12.5, CLAYD, True)]],
     align=PP_ALIGN.CENTER)
rect(s, Inches(0.55), Inches(5.95), Inches(12.23), Inches(0.95), RGBColor(0xEF,0xE7,0xDE), line=LINE)
text(s, Inches(0.8), Inches(6.08), Inches(11.8), Inches(0.8),
     [[("打法本质：", 12.5, CLAY, True),
       ("用「免费/开源的连接层」换「广泛分发」，再用「按量计费的智能层」捕获价值 —— 把自己做成 AI 应用生态绕不开的「水电煤」。", 12.5, INK, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_sp=1.05)

# ============================================================
# 5. 标杆案例
# ============================================================
s = new_slide("标杆案例：数据平台 × 编程赛道的双轮放量", "数据平台带来「治理数据上的企业 Agent」，编程赛道带来「最高 Token 消耗」", 4, TOTAL)

cases = [
    (BLUE, "Snowflake", "$200M · 多年期",
     ["触达 12,600+ 全球企业客户",
      "Claude 驱动 Cortex Code / Snowflake Intelligence",
      "联合 GTM：面向大型企业部署 AI Agent",
      "上架 Claude Marketplace"]),
    (GREEN, "Databricks", "5 年期战略合作",
     ["Claude 直接进入 Data Intelligence 平台",
      "经 Unity AI Gateway 治理与调用",
      "在「治理数据」上构建领域专用 Agent",
      "获评 2026 ISV AI 转型年度伙伴"]),
    (CLAY, "编程赛道", "最大收入引擎",
     ["Copilot / Cursor / Windsurf 全部接入 Claude",
      "份额：3% (25.04) → 18% (26.01)，9 月 6×",
      "单家企业级客户日均百万次 → 七位数年消耗",
      "Claude Code 自营：>50% 收入来自企业"]),
]
cx, cw, gap, cy, ch = Inches(0.55), Inches(3.97), Inches(0.18), Inches(1.62), Inches(3.7)
for i, (col, name, tag, items) in enumerate(cases):
    x = x_of(cx, cw, gap, i)
    rect(s, x, cy, cw, ch, CARD, line=LINE)
    rect(s, x, cy, cw, Inches(1.0), col)
    text(s, Emu(int(x)+Inches(0.22)), Emu(int(cy)+Inches(0.14)), Emu(int(cw)-Inches(0.4)), Inches(0.85),
         [[(name, 18, WHITE, True)], [(tag, 11.5, RGBColor(0xF2,0xF2,0xF2), True)]], sp_after=1)
    lines = [[("▪ ", 11, col, True), (it, 11.5, INK, False)] for it in items]
    text(s, Emu(int(x)+Inches(0.22)), Emu(int(cy)+Inches(1.18)), Emu(int(cw)-Inches(0.44)), Inches(2.4),
         lines, sp_after=10, line_sp=1.05)

text(s, Inches(0.55), Inches(5.55), Inches(12.23), Inches(0.95),
     [[("共性逻辑：", 12.5, CLAY, True),
       ("Anthropic 不直接卖给终端企业，而是「嫁接」到已握有数据与客户的平台上 —— 平台负责合规、数据与渠道，Claude 负责智能，", 12, INK, False)],
      [("            按 Token 用量分成。编程之所以成为收入引擎，正因其「长上下文 + 高输出 + 高频重复」，是天然的算力消耗黑洞。", 12, INK, False)]],
     sp_after=3, line_sp=1.1)

# ============================================================
# 6. MCP 标准战略
# ============================================================
s = new_slide("MCP：把「连接标准」做成生态护城河", "Model Context Protocol —— AI 时代的「USB-C」，用开源换主导权", 5, TOTAL)

# 左：时间线
tx = Inches(0.55)
text(s, tx, Inches(1.55), Inches(5.4), Inches(0.4), [[("演进路径", 15, CLAY, True)]])
timeline = [
    ("2024.11", "Anthropic 开源发布 MCP", CLAY),
    ("2025.03", "OpenAI 官方采纳（含 ChatGPT）", BLUE),
    ("2025.04", "Google DeepMind 宣布 Gemini 支持", GREEN),
    ("2025.12", "捐入 Linux 基金会 Agentic AI Foundation\n（联合 Block、OpenAI 共建）", PURP),
    ("2026.03", "全主流厂商接入，成事实标准", GOLD),
]
ty = Inches(2.05)
for date, ev, col in timeline:
    rect(s, tx, ty, Inches(0.16), Inches(0.72), col)
    text(s, Emu(int(tx)+Inches(0.32)), ty, Inches(5.1), Inches(0.78),
         [[(date, 12, col, True), ("   "+ev.split(chr(10))[0], 11.5, INK, True)]]
         + ([[(ev.split(chr(10))[1], 10, SUB, False)]] if "\n" in ev else []),
         sp_after=1, line_sp=1.0)
    ty = Emu(int(ty) + Inches(0.88))

# 右：采纳数据 + 战略意义
rx = Inches(6.5)
text(s, rx, Inches(1.55), Inches(6.3), Inches(0.4), [[("采纳规模（2026 Q1）", 15, GREEN, True)]])
mcp_stats = [("10,000+", "活跃公共 MCP 服务器"), ("97M", "Python/TS SDK 月下载量"), ("41%", "软件组织已投入生产")]
bx = rx
for i, (big, d) in enumerate(mcp_stats):
    x = x_of(rx, Inches(2.05), Inches(0.08), i)
    rect(s, x, Inches(2.0), Inches(2.0), Inches(1.1), CARD, line=LINE)
    text(s, x, Inches(2.12), Inches(2.0), Inches(0.9),
         [[(big, 24, GREEN, True)], [(d, 9.5, SUB, False)]], align=PP_ALIGN.CENTER, sp_after=2)

text(s, rx, Inches(3.45), Inches(6.3), Inches(0.4), [[("为何「送出去」反而是赢", 15, CLAY, True)]])
why = [
    ("定标准者得生态", "连接协议一旦通用，Claude 是协议的「原生公民」，调度与体验最优。"),
    ("捐给基金会＝去厂商化", "消除「被单一公司绑架」顾虑，加速对手与中立方共同采纳。"),
    ("做大蛋糕再分蛋糕", "MCP 让接入 SaaS/工具的成本趋零，Agent 用例爆发 → Token 消耗整体放大。"),
]
wy = Inches(3.95)
for h, d in why:
    rect(s, rx, wy, Inches(6.3), Inches(0.82), RGBColor(0xEF,0xE7,0xDE), line=LINE)
    text(s, Emu(int(rx)+Inches(0.18)), Emu(int(wy)+Inches(0.08)), Inches(6.0), Inches(0.7),
         [[(h, 12.5, CLAYD, True)], [(d, 10.5, INK, False)]], sp_after=1, line_sp=1.0)
    wy = Emu(int(wy) + Inches(0.92))

# ============================================================
# 7. 竞合悖论
# ============================================================
s = new_slide("竞合悖论：既是 SaaS 的供应商，也是颠覆者", "Anthropic 一边给 SaaS 供货，一边沿应用层向上攀升 —— 倒金字塔顶端的 100× 价值同样诱人", 6, TOTAL)

# 左右两栏对照
lx = Inches(0.55)
rect(s, lx, Inches(1.62), Inches(6.05), Inches(3.35), CARD, line=LINE)
rect(s, lx, Inches(1.62), Inches(6.05), Inches(0.62), GREEN)
text(s, Emu(int(lx)+Inches(0.22)), Inches(1.70), Inches(5.6), Inches(0.5),
     [[("作为「合作者」—— 成分品牌", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
coop = [
    "嵌入 Salesforce/ServiceNow/Snowflake 等核心系统",
    "MCP 成为 Agent 时代的「思考引擎层」",
    "增强而非取代记录系统（System of Record）",
    "让 SaaS 产品体验更好、留存更高",
]
text(s, Emu(int(lx)+Inches(0.28)), Inches(2.45), Inches(5.5), Inches(2.4),
     [[("✓ ", 12, GREEN, True), (c, 12, INK, False)] for c in coop], sp_after=10, line_sp=1.1)

rx2 = Inches(6.73)
rect(s, rx2, Inches(1.62), Inches(6.05), Inches(3.35), CARD, line=LINE)
rect(s, rx2, Inches(1.62), Inches(6.05), Inches(0.62), CLAYD)
text(s, Emu(int(rx2)+Inches(0.22)), Inches(1.70), Inches(5.6), Inches(0.5),
     [[("作为「颠覆者」—— 向上攀升", 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
comp = [
    "Claude Cowork 等让用户绕过 SaaS 自有界面",
    "按结果计费 → 侵蚀 SaaS「按席位」订阅模式",
    "法务/HR/财务垂直插件，直接做工作流",
    "Claude Code 自营，与编程类 SaaS 正面竞争",
]
text(s, Emu(int(rx2)+Inches(0.28)), Inches(2.45), Inches(5.5), Inches(2.4),
     [[("⚡ ", 12, CLAYD, True), (c, 12, INK, False)] for c in comp], sp_after=10, line_sp=1.1)

# 底部结论
rect(s, Inches(0.55), Inches(5.25), Inches(12.23), Inches(1.35), DARK)
text(s, Inches(0.85), Inches(5.42), Inches(11.7), Inches(1.1),
     [[("生态位判断：", 13.5, CLAY, True),
       ("Anthropic 正从「模型供应商」演化为「平台」。对 SaaS 厂商而言 —— 拥有专有数据、深度工作流与合规壁垒者，", 12.5, WHITE, False)],
      [("            被「赋能」（更强）；仅靠通用 AI 包装的薄应用，则可能被「商品化」（被取代）。这正是合作中必须看清的边界。", 12.5, WHITE, False)]],
     sp_after=3, line_sp=1.1)

# ============================================================
# 8. 核心方法论与启示
# ============================================================
s = new_slide("核心方法论：模型厂商「向上要价值」的四步法", "对行业智能化项目的启示 —— 价值不在模型本身，而在模型被嵌入业务的深度", 7, TOTAL)

steps = [
    (CLAY, "1", "成分化",
     "不直接卖模型，而是嵌入 SaaS 成为「智能成分」，用按量计费搭上应用层 100× 的价值车。"),
    (BLUE, "2", "借船分发",
     "嫁接握有数据与客户的平台（数据平台 / 行业 SaaS），让对方负责合规与渠道，自己专注智能。"),
    (GREEN, "3", "立标准",
     "开源 + 捐入中立基金会，把连接协议（MCP）做成事实标准，以「去厂商化」换生态主导权。"),
    (PURP, "4", "选点上攀",
     "在 Token 消耗最高、价值最厚的场景（编程、Agent）选择性下场自营，捕获应用层利润。"),
]
sx, sw, sgap, sy, sh2 = Inches(0.55), Inches(2.95), Inches(0.18), Inches(1.62), Inches(2.55)
for i, (col, num, name, d) in enumerate(steps):
    x = x_of(sx, sw, sgap, i)
    rect(s, x, sy, sw, sh2, CARD, line=LINE)
    rect(s, x, sy, sw, Inches(0.62), col)
    c = rect(s, Emu(int(x)+Inches(0.2)), Emu(int(sy)+Inches(0.78)), Inches(0.62), Inches(0.62), col, shape=MSO_SHAPE.OVAL)
    text(s, Emu(int(x)+Inches(0.2)), Emu(int(sy)+Inches(0.80)), Inches(0.62), Inches(0.58),
         [[(num, 22, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(x)+Inches(0.2)), Inches(1.70), Emu(int(sw)-Inches(0.4)), Inches(0.5),
         [[(name, 16, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(x)+Inches(0.22)), Emu(int(sy)+Inches(1.55)), Emu(int(sw)-Inches(0.42)), Inches(1.0),
         [[(d, 11, INK, False)]], line_sp=1.12)

# 启示条
rect(s, Inches(0.55), Inches(4.5), Inches(12.23), Inches(2.0), RGBColor(0xEF,0xE7,0xDE), line=LINE)
rect(s, Inches(0.55), Inches(4.5), Inches(0.14), Inches(2.0), CLAY)
text(s, Inches(0.9), Inches(4.65), Inches(11.6), Inches(0.5),
     [[("对行业智能化项目的三点启示", 16, CLAYD, True)]])
takes = [
    ("价值在「最后一公里」", "模型只是中间品；行业智能化的超额价值，来自把模型封装成可交付的业务结果。"),
    ("生态优先于单点产品", "成功项目不是「做一个 App」，而是嵌入既有数据/客户/合规体系，借生态放量。"),
    ("标准与数据是护城河", "谁掌握连接标准（如 MCP）与专有行业数据，谁就在倒金字塔顶端不被商品化。"),
]
ty = Inches(5.25)
for i, (h, d) in enumerate(takes):
    x = x_of(Inches(0.95), Inches(3.95), Inches(0.05), i)
    text(s, x, ty, Inches(3.85), Inches(1.1),
         [[("● ", 12, CLAY, True), (h, 12.5, INK, True)], [(d, 11, SUB, False)]],
         sp_after=3, line_sp=1.08)

text(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.3),
     [[("资料来源：Anthropic / Snowflake / Databricks 官方公告，Constellation Research，Sacra，Wikipedia(MCP)，行业报道（2024.11–2026.03）；部分为估算口径", 8, SUB, False)]])

out = "/Users/apple/Future_Thoughts/Anthropic生态分析.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
