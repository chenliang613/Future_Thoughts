# -*- coding: utf-8 -*-
"""Generate 微软AI治理分析.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
MS_BLUE   = RGBColor(0x00, 0x78, 0xD4)   # Microsoft blue
DEEP      = RGBColor(0x0B, 0x2E, 0x4F)   # deep navy
ACCENT    = RGBColor(0x50, 0xB0, 0x83)   # green accent
AMBER     = RGBColor(0xF2, 0x9B, 0x2E)   # amber
GREY      = RGBColor(0x5A, 0x63, 0x6E)   # body grey
LIGHT     = RGBColor(0xF2, 0xF6, 0xFB)   # light panel
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1B, 0x1B, 0x1B)
RED       = RGBColor(0xD0, 0x4A, 0x44)

FONT = "Microsoft YaHei"
FONT_L = "Microsoft YaHei Light"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_font(run, size, color, bold=False, font=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    # ensure east-asian font
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {})
        rpr.append(ea)
    ea.set('typeface', font)


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None, shadow=False, round_=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def para(tf, text, size, color, bold=False, font=FONT, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, level=0, first=False, line_spacing=None,
         italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    _set_font(r, size, color, bold, font, italic)
    return p


def page_chrome(slide, kicker, title, idx):
    """Standard content-slide header."""
    rect(slide, 0, 0, SW, Inches(1.18), WHITE)
    rect(slide, 0, 0, Inches(0.22), Inches(1.18), MS_BLUE)
    tb, tf = textbox(slide, Inches(0.55), Inches(0.16), Inches(11.6), Inches(0.95))
    para(tf, kicker, 12, MS_BLUE, bold=True, first=True, space_after=2)
    para(tf, title, 25, DEEP, bold=True, space_after=0)
    # footer
    fb, ff = textbox(slide, Inches(0.55), Inches(7.02), Inches(10), Inches(0.4))
    para(ff, "微软AI治理分析  ·  行业智能化产业发展研究", 9, GREY, first=True)
    nb, nf = textbox(slide, Inches(12.3), Inches(7.02), Inches(0.8), Inches(0.4))
    para(nf, str(idx), 10, GREY, align=PP_ALIGN.RIGHT, first=True)


def bullet(tf, text, size=14, color=GREY, bold=False, first=False, marker="▍",
           mcolor=MS_BLUE, space_after=9, sub=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = 1.05
    if sub:
        r0 = p.add_run(); r0.text = "– "; _set_font(r0, size, GREY, False)
    else:
        r0 = p.add_run(); r0.text = marker + "  "; _set_font(r0, size, mcolor, True)
    r = p.add_run(); r.text = text; _set_font(r, size, color, bold)
    return p


# =================================================================
# 1. COVER
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0, 0, SW, Inches(0.28), MS_BLUE)
rect(s, 0, SH - Inches(0.18), SW, Inches(0.18), ACCENT)
# decorative band
rect(s, Inches(8.6), 0, Inches(0.06), SH, RGBColor(0x14, 0x3A, 0x5E))
tb, tf = textbox(s, Inches(0.9), Inches(2.2), Inches(9.5), Inches(2.6))
para(tf, "AI 安全治理 · 战略分析", 15, ACCENT, bold=True, first=True, space_after=14)
para(tf, "微软 AI 治理分析", 46, WHITE, bold=True, space_after=8)
para(tf, "以 AI 安全与治理保障全球市场准入，并构建产品竞争壁垒", 19, RGBColor(0xC9, 0xDA, 0xEE))
tb2, tf2 = textbox(s, Inches(0.92), Inches(6.2), Inches(11), Inches(0.8))
para(tf2, "行业智能化项目和产业发展分析  |  研究报告  |  2026年6月", 13, RGBColor(0x9F, 0xB6, 0xCF), first=True)

# =================================================================
# 2. EXECUTIVE SUMMARY
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "EXECUTIVE SUMMARY", "执行摘要：治理是微软的“双重武器”", 2)
# central thesis band
rect(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.95), LIGHT, round_=True)
tb, tf = textbox(s, Inches(0.8), Inches(1.5), Inches(11.8), Inches(0.85), MSO_ANCHOR.MIDDLE)
para(tf, "核心论点：微软将 AI 安全与治理同时用作“合规防御盾”和“商业进攻矛”——",
     15, DEEP, bold=True, first=True, space_after=3)
para(tf, "对内满足各国监管以保住市场准入，对外把“可信”包装为企业级产品卖点，形成难以复制的信任护城河。",
     13.5, GREY)
# two pillars
cards = [
    (MS_BLUE, "防御：保障市场准入",
     ["以《负责任 AI 标准》对齐欧盟 AI 法案等全球法规",
      "分层合规：系统审查 + 政策更新 + 合同约束",
      "跨职能治理组织 + 33 份透明度说明",
      "责任共担模型，支持客户履行下游合规义务"]),
    (ACCENT, "进攻：提升竞争力",
     ["Trustworthy AI：安全、隐私、安全性三支柱",
      "把治理能力产品化（Content Safety / Foundry / Purview）",
      "客户版权承诺(CCC)等承诺降低客户采用风险",
      "以“企业级可信”锁定大客户与受监管行业"]),
]
cx = Inches(0.55)
cw = Inches(5.95)
for color, head, items in cards:
    rect(s, cx, Inches(2.7), cw, Inches(3.9), WHITE, line=RGBColor(0xE0,0xE6,0xED), round_=True)
    rect(s, cx, Inches(2.7), cw, Inches(0.7), color, round_=True)
    hb, hf = textbox(s, cx + Inches(0.3), Inches(2.78), cw - Inches(0.5), Inches(0.55), MSO_ANCHOR.MIDDLE)
    para(hf, head, 16, WHITE, bold=True, first=True)
    bb, bf = textbox(s, cx + Inches(0.32), Inches(3.6), cw - Inches(0.6), Inches(2.9))
    for i, it in enumerate(items):
        bullet(bf, it, 13, GREY, first=(i==0), mcolor=color, space_after=11)
    cx += cw + Inches(0.3)

# =================================================================
# 3. AGENDA
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "CONTENTS", "报告框架", 3)
items = [
    ("01", "全球 AI 监管格局与欧盟 AI 法案", "市场准入门槛正在抬高"),
    ("02", "微软 AI 治理的顶层框架", "六大原则与治理组织"),
    ("03", "第一部分 · 应对法案、保障市场准入", "分层合规与多法域策略"),
    ("04", "产品内嵌合规能力 & 客户责任共担", "把法规要求工程化"),
    ("05", "第二部分 · 以治理提升竞争力", "Trustworthy AI 与治理产品化"),
    ("06", "对行业智能化项目的启示", "可迁移的成功方法论"),
]
y = Inches(1.55)
for num, t, sub in items:
    rect(s, Inches(0.55), y, Inches(0.85), Inches(0.78), MS_BLUE, round_=True)
    nb, nf = textbox(s, Inches(0.55), y, Inches(0.85), Inches(0.78), MSO_ANCHOR.MIDDLE)
    para(nf, num, 20, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tb, tf = textbox(s, Inches(1.6), y, Inches(11), Inches(0.78), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.line_spacing = 1.0
    r = p.add_run(); r.text = t + "    "; _set_font(r, 16, DEEP, True)
    r2 = p.add_run(); r2.text = sub; _set_font(r2, 12.5, GREY, False)
    y += Inches(0.9)

# =================================================================
# 4. REGULATORY LANDSCAPE / EU AI ACT TIMELINE
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "背景 · 监管格局", "为何治理事关生死：监管正成为市场准入门槛", 4)
tb, tf = textbox(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.7))
para(tf, "欧盟 AI 法案是全球首部综合性 AI 法律，采用“风险分级 + 域外效力 + 高额罚款”，事实上成为全球基准。"
         "不合规即丧失欧盟市场准入——对微软这类全球平台厂商，治理已是经营底线。", 13.5, GREY, first=True)
# timeline
phases = [
    ("2024.08", "法案生效", "AI Act 正式进入立法生效阶段", MS_BLUE),
    ("2025.02", "禁止性条款", "社会评分等被禁用途生效；AI 素养义务启动", RED),
    ("2025.08", "通用模型(GPAI)", "通用目的 AI 模型义务 + 治理架构生效", AMBER),
    ("2026.08", "高风险系统", "高风险系统全面义务：风险管理、技术文档、人为监督", ACCENT),
]
n = len(phases)
x0 = Inches(0.7); span = Inches(11.9); cardw = Inches(2.75)
gap = (span - cardw * n)
step = (span - cardw) / (n - 1)
# baseline
rect(s, x0, Inches(2.55), span, Inches(0.06), RGBColor(0xD5,0xDD,0xE6))
for i, (date, head, desc, color) in enumerate(phases):
    cx = x0 + Emu(int(step * i))
    # node
    dot = rect(s, cx + cardw/2 - Inches(0.1), Inches(2.45), Inches(0.2), Inches(0.26), color, round_=True)
    rect(s, cx, Inches(2.95), cardw, Inches(0.55), color, round_=True)
    db, df = textbox(s, cx, Inches(2.97), cardw, Inches(0.5), MSO_ANCHOR.MIDDLE)
    para(df, date + "  " + head, 13, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    cb, cf = textbox(s, cx + Inches(0.05), Inches(3.6), cardw - Inches(0.1), Inches(1.5))
    para(cf, desc, 11.5, GREY, first=True, align=PP_ALIGN.CENTER, line_spacing=1.05)
# bottom callout
rect(s, Inches(0.55), Inches(5.45), Inches(12.2), Inches(1.35), LIGHT, round_=True)
tb, tf = textbox(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(1.2), MSO_ANCHOR.MIDDLE)
para(tf, "三重压力 → 治理成为准入前提", 14, DEEP, bold=True, first=True, space_after=4)
para(tf, "① 域外效力：只要服务欧盟用户即受约束   ② 罚款最高达全球营收 7%   "
         "③ 多法域叠加：美国行政令/各州法、中国生成式AI管理办法、英国/日本框架同步推进", 12.5, GREY)

# =================================================================
# 5. MICROSOFT GOVERNANCE FRAMEWORK
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "治理框架", "微软 AI 治理的顶层架构：原则 → 标准 → 组织", 5)
# six principles
tb, tf = textbox(s, Inches(0.55), Inches(1.3), Inches(6), Inches(0.4))
para(tf, "负责任 AI 六大原则（2018 起）", 14, DEEP, bold=True, first=True)
principles = [
    ("公平性", "Fairness"), ("可靠与安全", "Reliability & Safety"),
    ("隐私与安全", "Privacy & Security"), ("包容性", "Inclusiveness"),
    ("透明性", "Transparency"), ("问责制", "Accountability"),
]
px, py = Inches(0.55), Inches(1.85)
pw, ph = Inches(2.92), Inches(1.0)
for i, (cn, en) in enumerate(principles):
    col = i % 2; row = i // 2
    x = px + (pw + Inches(0.2)) * col
    y = py + (ph + Inches(0.2)) * row
    rect(s, x, y, pw, ph, LIGHT, line=RGBColor(0xDD,0xE5,0xEE), round_=True)
    rect(s, x, y, Inches(0.1), ph, MS_BLUE)
    b, f = textbox(s, x + Inches(0.25), y, pw - Inches(0.3), ph, MSO_ANCHOR.MIDDLE)
    para(f, cn, 14.5, DEEP, bold=True, first=True, space_after=1)
    para(f, en, 10.5, GREY)
# right column: governance machinery
rx = Inches(6.9)
rect(s, rx, Inches(1.3), Inches(5.85), Inches(5.35), DEEP, round_=True)
b, f = textbox(s, rx + Inches(0.35), Inches(1.5), Inches(5.2), Inches(5.0))
para(f, "把原则落地的治理机器", 15, ACCENT, bold=True, first=True, space_after=10)
gov = [
    ("负责任 AI 标准 (RAI Standard)", "把六大原则转化为可执行的工程要求与审查门槛"),
    ("负责任 AI 办公室 + RAI 委员会", "制定政策、裁定高风险案例、自上而下治理"),
    ("影响评估 + 红队对抗测试", "在研发全生命周期识别危害（PyRIT 框架）"),
    ("Trustworthy AI 倡议 (2024)", "整合安全 / 隐私 / 安全性，对齐安全未来计划(SFI)"),
    ("年度透明度报告", "对外披露治理实践，建立公众与监管信任"),
]
for i, (h, d) in enumerate(gov):
    p = f.add_paragraph(); p.space_after = Pt(4); p.space_before = Pt(2)
    r = p.add_run(); r.text = "● "; _set_font(r, 13, ACCENT, True)
    r2 = p.add_run(); r2.text = h; _set_font(r2, 13.5, WHITE, True)
    p2 = f.add_paragraph(); p2.space_after = Pt(9); p2.line_spacing = 1.0
    r3 = p2.add_run(); r3.text = "   " + d; _set_font(r3, 11.5, RGBColor(0xBF,0xD2,0xE6), False)

# =================================================================
# SECTION DIVIDER 1
# =================================================================
def divider(part, title, sub, idx):
    s = add_slide()
    rect(s, 0, 0, SW, SH, DEEP)
    rect(s, 0, 0, Inches(0.28), SH, MS_BLUE)
    rect(s, Inches(0.9), Inches(2.5), Inches(1.7), Inches(0.1), ACCENT)
    tb, tf = textbox(s, Inches(0.9), Inches(2.8), Inches(11), Inches(2.5))
    para(tf, part, 16, ACCENT, bold=True, first=True, space_after=14)
    para(tf, title, 40, WHITE, bold=True, space_after=10)
    para(tf, sub, 16, RGBColor(0xBF,0xD2,0xE6))
    nb, nf = textbox(s, Inches(12.3), Inches(7.02), Inches(0.8), Inches(0.4))
    para(nf, str(idx), 10, RGBColor(0x9F,0xB6,0xCF), align=PP_ALIGN.RIGHT, first=True)
    return s

divider("第一部分 · 防御", "应对各国法案，保障市场准入",
        "微软如何用治理体系把全球监管转化为可持续的市场通行证", 6)

# =================================================================
# 7. COMPLIANCE STRATEGY — 3 LAYERS
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "合规策略", "应对欧盟 AI 法案：分层合规“三板斧”", 7)
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.55))
para(tf, "面对 2025.02 禁止性条款生效，微软在欧盟 AI 办公室细则出台前即采取主动、分层的合规动作：", 13.5, GREY, first=True)
layers = [
    ("01 系统审查", MS_BLUE,
     ["设计筛查问卷，识别内部 AI 系统是否触及禁止用途",
      "对全部在售自有系统进行排查",
      "对被标记系统要求团队整改并复核"]),
    ("02 政策更新", AMBER,
     ["新增“受限用途”政策，禁止设计/部署禁止性应用",
      "把法规要求内化为内部研发红线",
      "与负责任 AI 标准联动执行"]),
    ("03 合同约束", ACCENT,
     ["更新客户协议与《生成式 AI 行为准则》",
      "明确禁止社会评分等违规用途",
      "把合规义务向客户与合作伙伴传导"]),
]
cx = Inches(0.55); cw = Inches(3.95)
for head, color, items in layers:
    rect(s, cx, Inches(2.05), cw, Inches(4.0), WHITE, line=RGBColor(0xE0,0xE6,0xED), round_=True)
    rect(s, cx, Inches(2.05), cw, Inches(0.78), color, round_=True)
    hb, hf = textbox(s, cx, Inches(2.1), cw, Inches(0.68), MSO_ANCHOR.MIDDLE)
    para(hf, head, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    bb, bf = textbox(s, cx + Inches(0.3), Inches(3.05), cw - Inches(0.55), Inches(2.9))
    for i, it in enumerate(items):
        bullet(bf, it, 13, GREY, first=(i==0), mcolor=color, space_after=12)
    cx += cw + Inches(0.27)
rect(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.55), LIGHT, round_=True)
b, f = textbox(s, Inches(0.8), Inches(6.27), Inches(11.7), Inches(0.5), MSO_ANCHOR.MIDDLE)
para(f, "关键特征：抢在监管细则之前“proactive compliance”——把合规节奏掌握在自己手里，降低突击整改风险。",
     12.5, DEEP, bold=True, first=True)

# =================================================================
# 8. PRODUCTS EMBED COMPLIANCE
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "产品工程化", "把法规要求“工程化”：合规能力内嵌进产品", 8)
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.55))
para(tf, "微软不止于政策文件，而是将透明度、安全、可评估等法规要求转化为可调用的产品与工具：", 13.5, GREY, first=True)
tools = [
    ("透明度说明 Transparency Notes", "自 2019 年发布 33+ 份，覆盖各 AI 工具的能力与局限，直接回应法案“透明度”义务", MS_BLUE),
    ("Azure AI Content Safety", "内容安全过滤，拦截有害/不当输出，缓解危害风险", RED),
    ("Azure AI Foundry 评估", "对生成式 AI 做系统化评估与监控，支撑高风险系统的测试义务", ACCENT),
    ("PyRIT 开源红队框架", "对抗性测试识别危害，把安全测试标准化、可复用", AMBER),
    ("Microsoft Purview / Compliance Manager", "多法规合规追踪与数据治理，统一管理跨法域义务", DEEP),
    ("内置防护：防注入 / 防版权侵权", "AI 服务自带基础安全控制（Secure by Default）", MS_BLUE),
]
px, py = Inches(0.55), Inches(2.0)
cw, ch = Inches(5.98), Inches(1.42)
for i, (h, d, color) in enumerate(tools):
    col = i % 2; row = i // 2
    x = px + (cw + Inches(0.24)) * col
    y = py + (ch + Inches(0.18)) * row
    rect(s, x, y, cw, ch, LIGHT, line=RGBColor(0xDD,0xE5,0xEE), round_=True)
    rect(s, x, y, Inches(0.12), ch, color)
    b, f = textbox(s, x + Inches(0.32), y + Inches(0.12), cw - Inches(0.5), ch - Inches(0.2), MSO_ANCHOR.MIDDLE)
    para(f, h, 14, DEEP, bold=True, first=True, space_after=3)
    para(f, d, 11.5, GREY, line_spacing=1.02)

# =================================================================
# 9. SHARED RESPONSIBILITY / CUSTOMER COMMITMENTS
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "责任共担", "不只自己合规：帮客户合规的“责任共担”模型", 9)
# left
tb, tf = textbox(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(0.5))
para(tf, "供应商—部署者的责任共担", 15, DEEP, bold=True, first=True)
b, f = textbox(s, Inches(0.55), Inches(1.95), Inches(6.0), Inches(4.6))
pts = [
    "微软作为 AI 工具/组件提供方，须支持下游“受监管行为者”（企业客户）",
    "当客户把微软工具集成进高风险系统时，微软通过知识共享、文档、工具承接其部分合规负担",
    "AI 共担责任模型：明确划分平台方与客户各自的安全/合规边界",
    "2024.06 发布 AI 客户承诺(AI Customer Commitments)：陪伴客户走负责任 AI 旅程",
]
for i, t in enumerate(pts):
    bullet(f, t, 13.5, GREY, first=(i==0), space_after=14)
# right highlight: copyright commitment
rx = Inches(6.95)
rect(s, rx, Inches(1.4), Inches(5.8), Inches(5.1), DEEP, round_=True)
b, f = textbox(s, rx + Inches(0.4), Inches(1.7), Inches(5.1), Inches(4.6))
para(f, "招牌承诺：客户版权承诺 (CCC)", 16, ACCENT, bold=True, first=True, space_after=10)
para(f, "写入微软产品条款的法律义务：若客户因使用 Copilot/Azure OpenAI 的输出内容遭第三方知识产权索赔，"
        "由微软出面抗辩并承担相应赔偿。", 13.5, WHITE, space_after=14, line_spacing=1.1)
para(f, "战略意义", 13.5, ACCENT, bold=True, space_after=6)
for t in ["把客户最担心的法律风险“兜底”，扫清采用顾虑",
          "前提是客户使用内置的版权防护缓解措施 → 反向推动安全功能采用",
          "用承诺把“信任”做成可签约、可背书的商业条款"]:
    p = f.add_paragraph(); p.space_after = Pt(8); p.line_spacing = 1.05
    r = p.add_run(); r.text = "✓ "; _set_font(r, 13, ACCENT, True)
    r2 = p.add_run(); r2.text = t; _set_font(r2, 12.5, RGBColor(0xCF,0xDD,0xEC))

# =================================================================
# 10. MULTI-JURISDICTION
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "多法域 & 行业自律", "超越欧盟：多法域协同与规则共建", 10)
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.5))
para(tf, "微软既被动适配各法域差异，也主动参与规则制定，争取把自身实践写进行业标准——“做规则的人”：", 13.5, GREY, first=True)
rows = [
    ("欧盟 GPAI 行为准则", "与 OpenAI、谷歌、Anthropic 等 26 家机构共同签署通用目的 AI 模型行为准则", MS_BLUE),
    ("AI Pact 自愿承诺", "签署欧盟 AI Pact 三项核心自愿承诺，提前对齐法案精神", ACCENT),
    ("标准共建 CEN / CENELEC", "参与欧洲技术标准制定，把工程实践转化为行业通用规范", AMBER),
    ("与欧盟 AI 办公室对接", "就实施细则与口径保持沟通，降低解释不确定性", DEEP),
    ("跨法域适配", "对齐美国行政令/州法、英国/日本框架、中国生成式AI管理办法等多套规则", RED),
]
y = Inches(2.0)
for h, d, color in rows:
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.82), LIGHT, line=RGBColor(0xDD,0xE5,0xEE), round_=True)
    rect(s, Inches(0.55), y, Inches(0.14), Inches(0.82), color)
    b, f = textbox(s, Inches(0.95), y, Inches(4.2), Inches(0.82), MSO_ANCHOR.MIDDLE)
    para(f, h, 14.5, DEEP, bold=True, first=True)
    b2, f2 = textbox(s, Inches(5.2), y, Inches(7.4), Inches(0.82), MSO_ANCHOR.MIDDLE)
    para(f2, d, 12.5, GREY, first=True, line_spacing=1.0)
    y += Inches(0.92)

# =================================================================
# SECTION DIVIDER 2
# =================================================================
divider("第二部分 · 进攻", "以治理手段提升产品竞争力",
        "把“可信”从成本中心变成增长引擎与差异化壁垒", 11)

# =================================================================
# 12. TRUSTWORTHY AI THREE PILLARS
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "Trustworthy AI", "竞争力基座：Trustworthy AI 三支柱", 12)
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.5))
para(tf, "2024 年 9 月，微软在 AI Tour 上推出 Trustworthy AI，把分散的能力整合为对外统一的“可信”品牌叙事：", 13.5, GREY, first=True)
pillars = [
    ("安全 Security", MS_BLUE,
     ["公司头号优先级，对齐安全未来计划(SFI)",
      "三原则：默认安全 / 设计安全 / 运营安全",
      "AI 服务自带防注入、防版权侵权等控制"]),
    ("安全性 Safety", ACCENT,
     ["以 2018 负责任 AI 原则为基础",
      "构建—测试—监控，规避有害内容/偏见/滥用",
      "影响评估 + 红队测试贯穿全生命周期"]),
    ("隐私 Privacy", AMBER,
     ["为企业提供细粒度数据控制",
      "如 M365 Copilot 可开关 Web 搜索等功能",
      "Copilot 在企业级安全/合规边界内运行"]),
]
cx = Inches(0.55); cw = Inches(3.95)
for head, color, items in pillars:
    rect(s, cx, Inches(2.05), cw, Inches(4.3), WHITE, line=RGBColor(0xE0,0xE6,0xED), round_=True)
    rect(s, cx, Inches(2.05), cw, Inches(0.85), color, round_=True)
    hb, hf = textbox(s, cx, Inches(2.1), cw, Inches(0.75), MSO_ANCHOR.MIDDLE)
    para(hf, head, 17, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    bb, bf = textbox(s, cx + Inches(0.3), Inches(3.15), cw - Inches(0.55), Inches(3.0))
    for i, it in enumerate(items):
        bullet(bf, it, 13, GREY, first=(i==0), mcolor=color, space_after=13)
    cx += cw + Inches(0.27)
rect(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4), WHITE)
b, f = textbox(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.45), MSO_ANCHOR.MIDDLE)
para(f, "“释放人类潜能始于信任” —— 把治理叙事提升为 CEO 级的市场定位。", 12.5, DEEP, bold=True, first=True, align=PP_ALIGN.CENTER)

# =================================================================
# 13. GOVERNANCE AS PRODUCT
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "治理产品化", "把治理变成卖点：合规能力即产品功能", 13)
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.5))
para(tf, "微软的差异化在于：竞争对手把安全治理当成本，微软把它做成可计费、可演示、可采购的能力。", 13.5, GREY, first=True)
prod = [
    ("治理工具自身即收入", "Content Safety、Foundry 评估、Purview 等作为付费/增值能力出售，治理投入反哺营收"),
    ("降低客户采用门槛", "CCC 版权兜底 + 透明度说明 + 共担责任，让受监管行业客户敢用、能合规地用"),
    ("一站式合规平台", "Azure 作为统一开发平台，提供企业级隐私/安全/合规，客户无需自建治理栈"),
    ("加速企业落地速度", "现成的评估、红队、内容过滤能力缩短客户的合规上线周期，形成速度优势"),
]
px, py = Inches(0.55), Inches(2.0)
cw, ch = Inches(6.0), Inches(2.2)
for i, (h, d) in enumerate(prod):
    col = i % 2; row = i // 2
    x = px + (cw + Inches(0.22)) * col
    y = py + (ch + Inches(0.22)) * row
    rect(s, x, y, cw, ch, LIGHT, line=RGBColor(0xDD,0xE5,0xEE), round_=True)
    rect(s, x, y, cw, Inches(0.12), MS_BLUE)
    b, f = textbox(s, x + Inches(0.35), y + Inches(0.3), cw - Inches(0.6), ch - Inches(0.5), MSO_ANCHOR.MIDDLE)
    para(f, h, 15.5, DEEP, bold=True, first=True, space_after=7)
    para(f, d, 13, GREY, line_spacing=1.1)

# =================================================================
# 14. TRUST AS MOAT — value loop
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "信任护城河", "商业闭环：信任如何转化为竞争壁垒", 14)
flow = [
    ("治理投入", "原则/标准/组织/工具", MS_BLUE),
    ("合规通行证", "保住欧盟等全球市场", AMBER),
    ("企业级信任", "受监管行业敢采用", ACCENT),
    ("锁定大客户", "高切换成本+续约", DEEP),
    ("数据与规模", "反哺治理与模型", MS_BLUE),
]
n = len(flow); x0 = Inches(0.6); bw = Inches(2.05); gap = Inches(0.42)
y = Inches(1.95)
for i, (h, d, color) in enumerate(flow):
    x = x0 + (bw + gap) * i
    rect(s, x, y, bw, Inches(1.5), color, round_=True)
    b, f = textbox(s, x + Inches(0.1), y, bw - Inches(0.2), Inches(1.5), MSO_ANCHOR.MIDDLE)
    para(f, h, 15, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=4)
    para(f, d, 11, RGBColor(0xE6,0xEE,0xF6), align=PP_ALIGN.CENTER, line_spacing=1.0)
    if i < n - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.04), y + Inches(0.55), gap - Inches(0.08), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(0xB8,0xC6,0xD6); ar.line.fill.background()
        ar.shadow.inherit = False
# loop-back note
rect(s, Inches(0.6), Inches(3.9), Inches(12.13), Inches(0.5), LIGHT, round_=True)
b, f = textbox(s, Inches(0.6), Inches(3.88), Inches(12.13), Inches(0.5), MSO_ANCHOR.MIDDLE)
para(f, "↺ 规模与数据再投入治理与模型，循环自我强化 —— 形成对手难以追赶的“信任飞轮”。",
     12.5, DEEP, bold=True, align=PP_ALIGN.CENTER, first=True)
# differentiators
tb, tf = textbox(s, Inches(0.6), Inches(4.6), Inches(12.13), Inches(0.4))
para(tf, "三大差异化结果", 14, DEEP, bold=True, first=True)
diffs = [
    ("先发的监管资本", "提前合规 + 参与立法，把规则话语权变成进入壁垒"),
    ("可信品牌溢价", "在 B2B 与受监管行业，“安全可信”是采购决策的硬指标"),
    ("生态绑定", "Azure 一体化治理栈提高迁移成本，强化平台黏性"),
]
cx = Inches(0.6); cw = Inches(3.92)
for h, d in diffs:
    rect(s, cx, Inches(5.1), cw, Inches(1.55), WHITE, line=RGBColor(0xDD,0xE5,0xEE), round_=True)
    rect(s, cx, Inches(5.1), Inches(0.1), Inches(1.55), ACCENT)
    b, f = textbox(s, cx + Inches(0.28), Inches(5.22), cw - Inches(0.45), Inches(1.3), MSO_ANCHOR.MIDDLE)
    para(f, h, 14, DEEP, bold=True, first=True, space_after=5)
    para(f, d, 12, GREY, line_spacing=1.05)
    cx += cw + Inches(0.18)

# =================================================================
# 15. IMPLICATIONS FOR 行业智能化
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
page_chrome(s, "启示与方法论", "对行业智能化项目的可迁移启示", 15)
tb, tf = textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.5))
para(tf, "微软的实践印证：在“倒金字塔”产业中，治理与信任正是行业应用层创造 100 倍价值的关键路径之一。", 13.5, GREY, first=True)
less = [
    ("治理前置，而非补救", "在项目设计期就嵌入安全/合规要求，把监管节奏掌握在自己手里"),
    ("把合规工程化", "将法规转化为可复用的工具与流程（评估、红队、内容过滤），而非一次性文档"),
    ("责任共担生态", "平台方为行业客户兜底关键风险，降低落地门槛、加速规模化"),
    ("信任即差异化", "在受监管行业，可信是采购硬指标——把治理做成卖点而非成本"),
    ("参与规则制定", "主动对接监管、共建标准，把自身实践变成行业门槛"),
    ("信任飞轮", "合规→采用→数据规模→再投入，构筑应用层的可持续护城河"),
]
px, py = Inches(0.55), Inches(2.0)
cw, ch = Inches(6.0), Inches(1.45)
for i, (h, d) in enumerate(less):
    col = i % 2; row = i // 2
    x = px + (cw + Inches(0.22)) * col
    y = py + (ch + Inches(0.16)) * row
    rect(s, x, y, cw, ch, LIGHT, line=RGBColor(0xDD,0xE5,0xEE), round_=True)
    nb = rect(s, x + Inches(0.2), y + Inches(0.28), Inches(0.5), Inches(0.5), MS_BLUE, round_=True)
    b0, f0 = textbox(s, x + Inches(0.2), y + Inches(0.28), Inches(0.5), Inches(0.5), MSO_ANCHOR.MIDDLE)
    para(f0, str(i+1), 16, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    b, f = textbox(s, x + Inches(0.85), y + Inches(0.12), cw - Inches(1.05), ch - Inches(0.2), MSO_ANCHOR.MIDDLE)
    para(f, h, 14, DEEP, bold=True, first=True, space_after=3)
    para(f, d, 11.5, GREY, line_spacing=1.02)

# =================================================================
# 16. CLOSING
# =================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, DEEP)
rect(s, 0, 0, SW, Inches(0.22), MS_BLUE)
rect(s, 0, SH - Inches(0.16), SW, Inches(0.16), ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.2))
para(tf, "结论", 16, ACCENT, bold=True, first=True, space_after=10)
para(tf, "治理不是 AI 的刹车，而是微软的方向盘与油门", 30, WHITE, bold=True)
b, f = textbox(s, Inches(0.92), Inches(3.1), Inches(11.4), Inches(3.0))
concl = [
    "防御端：以负责任 AI 标准对齐欧盟 AI 法案等全球法规，用“系统审查+政策+合同”分层合规保住市场准入。",
    "进攻端：以 Trustworthy AI 把安全/隐私/安全性整合为品牌，把治理能力产品化为可计费的竞争优势。",
    "护城河：合规→信任→锁定大客户→数据规模→再投入，形成自我强化的“信任飞轮”。",
    "对行业智能化：治理前置、合规工程化、责任共担、信任即差异化，是应用层创造倍增价值的关键路径。",
]
for i, t in enumerate(concl):
    p = f.add_paragraph(); p.space_after = Pt(15); p.line_spacing = 1.15
    r = p.add_run(); r.text = "▍ "; _set_font(r, 15, ACCENT, True)
    r2 = p.add_run(); r2.text = t; _set_font(r2, 15, RGBColor(0xDD,0xE7,0xF1))

prs.save("/Users/apple/Future_Thoughts/微软AI治理分析.pptx")
print("saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
print("OK")
