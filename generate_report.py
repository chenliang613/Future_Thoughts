#!/usr/bin/env python3
"""Generate Mythos and Glasswing Analysis Report PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MEDIUM = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_paragraph(tf, text, size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_bullet(tf, text, size=14, color=WHITE, level=0, bold=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

# ========== SLIDE 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Title accent line
add_shape(slide, Inches(1), Inches(2.5), Inches(0.08), Inches(2), ACCENT_BLUE)

tb = add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2))
set_text(tb.text_frame, "Claude Mythos & Project Glasswing", size=40, color=WHITE, bold=True)
tb.text_frame.word_wrap = True

tb2 = add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.8))
set_text(tb2.text_frame, "AI网络安全新范式 — 深度分析报告", size=24, color=ACCENT_BLUE)

tb3 = add_text_box(slide, Inches(1.5), Inches(5.0), Inches(8), Inches(0.5))
set_text(tb3.text_frame, "2026年4月  |  Anthropic前沿AI安全能力与产业协作分析", size=16, color=LIGHT_GRAY)

# ========== SLIDE 2: Executive Summary ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "执行摘要", size=32, color=WHITE, bold=True)
# underline
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), ACCENT_BLUE)

# Three summary cards
cards = [
    ("Claude Mythos Preview", "Anthropic迄今最强大的AI模型，在网络安全\n任务上实现了质的飞跃，能自主发现并利用\n主流操作系统和浏览器中的零日漏洞", ACCENT_BLUE),
    ("Project Glasswing", "一项协作性网络安全计划，将Mythos的\n能力限制在防御性用途，联合50+科技\n组织共同加固关键基础设施安全", ACCENT_ORANGE),
    ("产业影响", "标志着AI从\"辅助安全工具\"向\"自主安全\n研究员\"的跨越式转变，重塑攻防对抗\n格局，引发监管与伦理深度讨论", ACCENT_RED),
]

for i, (title, desc, accent) in enumerate(cards):
    left = Inches(0.8 + i * 4.1)
    card = add_shape(slide, left, Inches(1.6), Inches(3.7), Inches(4.5), BG_MEDIUM, accent)
    # Card title
    tb = add_text_box(slide, left + Inches(0.3), Inches(1.8), Inches(3.1), Inches(0.6))
    set_text(tb.text_frame, title, size=20, color=accent, bold=True)
    # Card desc
    tb2 = add_text_box(slide, left + Inches(0.3), Inches(2.6), Inches(3.1), Inches(3.2))
    tb2.text_frame.word_wrap = True
    set_text(tb2.text_frame, desc, size=15, color=LIGHT_GRAY)

# Key number strip
strip = add_shape(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8), RGBColor(0x0D, 0x15, 0x2D))
nums = [("数千个", "零日漏洞发现"), ("$1亿", "使用额度投入"), ("50+", "合作伙伴"), ("83.1%", "CyberGym得分")]
for i, (num, label) in enumerate(nums):
    x = Inches(1.2 + i * 3.0)
    tb = add_text_box(slide, x, Inches(6.3), Inches(2.5), Inches(0.45))
    set_text(tb.text_frame, num, size=22, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = add_text_box(slide, x, Inches(6.7), Inches(2.5), Inches(0.35))
    set_text(tb2.text_frame, label, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ========== SLIDE 3: What is Claude Mythos ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "什么是 Claude Mythos Preview?", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.05), ACCENT_BLUE)

# Left column - Definition
tb = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "模型定义", size=22, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=6)
bullets = [
    "Anthropic迄今构建的最强大AI模型",
    "通用语言模型，但在网络安全任务上表现惊人",
    "能力水平与Claude Opus 4.6相近，但安全领域远超",
    "代表了AI能力的\"阶跃式变化(Step Change)\"",
]
for b in bullets:
    add_bullet(tb.text_frame, "▸ " + b, size=15, color=LIGHT_GRAY)

add_paragraph(tb.text_frame, "", size=10)
add_paragraph(tb.text_frame, "核心安全能力", size=22, color=ACCENT_ORANGE, bold=True)
add_paragraph(tb.text_frame, "", size=6)
sec_bullets = [
    "自主发现零日漏洞（所有主流OS和浏览器）",
    "自主构建完整漏洞利用链（ROP链、堆喷射等）",
    "从剥离的二进制文件反向工程重建源代码",
    "发现的漏洞有些存活了17-27年之久",
    "99%以上发现的漏洞尚未被修补",
]
for b in sec_bullets:
    add_bullet(tb.text_frame, "▸ " + b, size=15, color=LIGHT_GRAY)

# Right column - Key case
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), BG_MEDIUM, ACCENT_BLUE)
tb = add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5.1), Inches(5))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "标志性发现案例", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=8)

cases = [
    ("FreeBSD NFS RCE (CVE-2026-4747)", "存活17年的远程代码执行漏洞，完全自主\n发现并利用，无需人工干预，可完全控制服务器"),
    ("OpenBSD TCP SACK漏洞", "存活27年的实现缺陷，可远程崩溃系统"),
    ("FFmpeg H.264编解码器漏洞", "存活16年，曾经历500万次自动化测试\n均未发现"),
    ("浏览器JIT堆喷射", "链接4个漏洞，逃逸渲染器和OS沙箱"),
    ("Linux内核特权提升", "多个漏洞链式组合实现权限提升"),
]
for title, desc in cases:
    add_paragraph(tb.text_frame, title, size=15, color=ACCENT_ORANGE, bold=True, space_before=Pt(12))
    add_paragraph(tb.text_frame, desc, size=13, color=LIGHT_GRAY, space_before=Pt(2))

# ========== SLIDE 4: Benchmarks ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "性能基准对比: Mythos vs Opus 4.6", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(4), Inches(0.05), ACCENT_BLUE)

benchmarks = [
    ("CyberGym", "83.1%", "66.6%", "+24.8%"),
    ("SWE-bench Verified", "93.9%", "80.8%", "+16.2%"),
    ("SWE-bench Pro", "77.8%", "53.4%", "+45.7%"),
    ("SWE-bench Multilingual", "87.3%", "77.8%", "+12.2%"),
    ("Terminal-Bench 2.0", "82.0%", "65.4%", "+25.4%"),
]

# Table header
header_y = Inches(1.6)
cols = [Inches(0.8), Inches(4.5), Inches(6.8), Inches(9.1), Inches(11.0)]
headers = ["基准测试", "Mythos Preview", "Opus 4.6", "提升幅度"]
header_bg = add_shape(slide, Inches(0.6), header_y, Inches(12.1), Inches(0.6), ACCENT_BLUE)

for j, (col_x, h_text) in enumerate(zip(cols, headers)):
    tb = add_text_box(slide, col_x, header_y + Inches(0.08), Inches(2), Inches(0.45))
    set_text(tb.text_frame, h_text, size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

for i, (name, mythos, opus, gain) in enumerate(benchmarks):
    row_y = Inches(2.4 + i * 0.85)
    bg_color = BG_MEDIUM if i % 2 == 0 else BG_DARK
    add_shape(slide, Inches(0.6), row_y, Inches(12.1), Inches(0.7), bg_color)

    vals = [name, mythos, opus, gain]
    colors = [WHITE, ACCENT_GREEN, LIGHT_GRAY, ACCENT_BLUE]
    for j, (col_x, val, clr) in enumerate(zip(cols, vals, colors)):
        tb = add_text_box(slide, col_x, row_y + Inches(0.1), Inches(2), Inches(0.5))
        set_text(tb.text_frame, val, size=16 if j == 0 else 18, color=clr, bold=(j > 0), alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# Firefox exploit comparison
card = add_shape(slide, Inches(0.8), Inches(6.0), Inches(5.5), Inches(1.1), BG_MEDIUM, ACCENT_ORANGE)
tb = add_text_box(slide, Inches(1.0), Inches(6.1), Inches(5.1), Inches(0.9))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "Firefox JS引擎漏洞利用: Mythos 181次 vs Opus 2次", size=15, color=ACCENT_ORANGE, bold=True)
add_paragraph(tb.text_frame, "成功率提升约90倍", size=13, color=LIGHT_GRAY)

# OSS-Fuzz
card = add_shape(slide, Inches(6.8), Inches(6.0), Inches(5.9), Inches(1.1), BG_MEDIUM, ACCENT_RED)
tb = add_text_box(slide, Inches(7.0), Inches(6.1), Inches(5.5), Inches(0.9))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "OSS-Fuzz: 595次崩溃 + 10次完整控制流劫持", size=15, color=ACCENT_RED, bold=True)
add_paragraph(tb.text_frame, "Opus 4.6仅实现单次三级崩溃", size=13, color=LIGHT_GRAY)

# ========== SLIDE 5: What is Project Glasswing ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "什么是 Project Glasswing?", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.05), ACCENT_ORANGE)

# Mission statement card
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.4), BG_MEDIUM, ACCENT_ORANGE)
tb = add_text_box(slide, Inches(1.2), Inches(1.7), Inches(11), Inches(1))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "核心使命", size=20, color=ACCENT_ORANGE, bold=True)
add_paragraph(tb.text_frame, "利用Claude Mythos Preview保护全球最关键的软件基础设施，在攻击者利用AI之前建立防御优势，\n同时为行业应对AI驱动的网络安全新时代做好准备。", size=15, color=LIGHT_GRAY)

# Why needed
tb = add_text_box(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(4))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "为什么需要Glasswing?", size=22, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=6)
reasons = [
    "Mythos的安全能力太强大，不适合公开发布",
    "恶意使用风险: 网络犯罪分子和间谍可利用其攻击",
    "漏洞发现和利用之间的窗口期已经崩塌",
    "传统安全缓解措施依赖的\"摩擦力\"面对AI大幅减弱",
    "N-day漏洞利用速度被大幅加速",
    "需要在攻击者之前修复关键漏洞",
]
for r in reasons:
    add_bullet(tb.text_frame, "▸ " + r, size=14, color=LIGHT_GRAY)

# Investment
tb = add_text_box(slide, Inches(6.8), Inches(3.2), Inches(5.7), Inches(4))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "资金投入", size=22, color=ACCENT_GREEN, bold=True)
add_paragraph(tb.text_frame, "", size=6)
investments = [
    ("$1亿", "Claude Mythos Preview使用额度"),
    ("$250万", "Alpha-Omega & OpenSSF (Linux基金会)"),
    ("$150万", "Apache软件基金会"),
]
for amount, target in investments:
    add_paragraph(tb.text_frame, amount, size=20, color=ACCENT_GREEN, bold=True, space_before=Pt(14))
    add_paragraph(tb.text_frame, target, size=14, color=LIGHT_GRAY, space_before=Pt(2))

add_paragraph(tb.text_frame, "", size=10)
add_paragraph(tb.text_frame, "访问渠道", size=22, color=ACCENT_ORANGE, bold=True)
add_paragraph(tb.text_frame, "", size=4)
channels = ["Claude API", "Amazon Bedrock", "Google Cloud Vertex AI", "Microsoft Foundry"]
for c in channels:
    add_bullet(tb.text_frame, "▸ " + c, size=14, color=LIGHT_GRAY)

# ========== SLIDE 6: Glasswing Partners ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "Glasswing 合作伙伴生态", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), ACCENT_ORANGE)

# 12 founding partners
partners = [
    ("Amazon Web Services", "云计算"), ("Apple", "消费电子"),
    ("Broadcom", "半导体"), ("Cisco", "网络设备"),
    ("CrowdStrike", "端点安全"), ("Google", "云/搜索"),
    ("JPMorgan Chase", "金融"), ("Linux Foundation", "开源"),
    ("Microsoft", "云/操作系统"), ("NVIDIA", "GPU/AI"),
    ("Palo Alto Networks", "网络安全"), ("Anthropic", "AI安全"),
]

tb_label = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5))
set_text(tb_label.text_frame, "12家创始合作伙伴", size=20, color=ACCENT_ORANGE, bold=True)

for i, (name, domain) in enumerate(partners):
    row = i // 4
    col = i % 4
    left = Inches(0.8 + col * 3.1)
    top = Inches(2.1 + row * 1.4)
    card = add_shape(slide, left, top, Inches(2.8), Inches(1.1), BG_MEDIUM, ACCENT_BLUE)
    tb = add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(2.4), Inches(0.5))
    set_text(tb.text_frame, name, size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(2.4), Inches(0.4))
    set_text(tb2.text_frame, domain, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Additional info
card = add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0), BG_MEDIUM, ACCENT_GREEN)
tb = add_text_box(slide, Inches(1.2), Inches(6.15), Inches(11), Inches(0.7))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "另有40+组织获得工具和资源访问权限  |  覆盖云计算、操作系统、网络安全、金融、开源等关键领域", size=16, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ========== SLIDE 7: Risks & Controversy ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "风险、争议与监管动态", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.05), ACCENT_RED)

# Risks
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), BG_MEDIUM, ACCENT_RED)
tb = add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(2.2))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "安全风险", size=20, color=ACCENT_RED, bold=True)
risks = [
    "99%+漏洞尚未修补，若泄露将造成灾难",
    "利用成本极低: 完整漏洞利用链<$1,000",
    "1000次自动化攻击尝试成本<$20,000",
    "传统安全\"摩擦力\"防线面对AI大幅弱化",
    "N-day利用窗口从数周缩短至数小时",
]
for r in risks:
    add_bullet(tb.text_frame, "⚠ " + r, size=13, color=LIGHT_GRAY)

# Controversy
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.5), BG_MEDIUM, ACCENT_ORANGE)
tb = add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(2.2))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "争议声音", size=20, color=ACCENT_ORANGE, bold=True)
controversies = [
    "Tom's Hardware: \"不是超级黑客，而是销售策略\"",
    "\"数千个\"严重漏洞仅有198个人工复核",
    "89%准确率意味着仍有11%误报",
    "模型泄露事件引发信息安全担忧",
    "公开发布延迟反映内部对风险的审慎评估",
]
for c in controversies:
    add_bullet(tb.text_frame, "▸ " + c, size=13, color=LIGHT_GRAY)

# Regulatory
card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8), BG_MEDIUM, ACCENT_BLUE)
tb = add_text_box(slide, Inches(1.1), Inches(4.5), Inches(11), Inches(2.5))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "监管与政府互动", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=6)

reg_items = [
    "与CISA（网络安全和基础设施安全局）持续讨论Mythos的网络能力",
    "与CAISI（AI标准与创新中心）就AI安全标准进行协商",
    "美联储主席Powell和财政部长Bessent与主要银行CEO讨论Mythos的网络安全影响",
    "FBI报告显示网络犯罪损失接近$210亿，强化了AI安全工具的紧迫性",
    "Anthropic主动延迟Mythos公开发布，优先确保安全防护措施到位",
    "呼吁政府参与维护技术领先地位和管控国家安全风险",
]
for item in reg_items:
    add_bullet(tb.text_frame, "▸ " + item, size=14, color=LIGHT_GRAY)

# ========== SLIDE 8: Industry Impact Analysis ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "产业影响分析", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.05), ACCENT_BLUE)

# Impact areas
areas = [
    ("网络安全行业", ACCENT_RED, [
        "漏洞发现从人工驱动转向AI自主",
        "安全研究员角色重新定义",
        "红队/蓝队对抗进入AI时代",
        "安全厂商需快速整合AI能力",
    ]),
    ("软件开发", ACCENT_BLUE, [
        "代码审计效率数量级提升",
        "开源项目安全性显著改善",
        "\"安全左移\"获得强力AI工具支撑",
        "开发流程需内嵌AI安全检查",
    ]),
    ("AI产业", ACCENT_GREEN, [
        "证明AI在垂直领域的巨大价值",
        "\"负责任发布\"成为行业新范式",
        "安全成为AI差异化竞争新维度",
        "推动AI监管框架加速建立",
    ]),
    ("金融/关基行业", ACCENT_ORANGE, [
        "关键基础设施防御能力大幅提升",
        "金融系统安全审计方式变革",
        "合规要求可能纳入AI安全工具",
        "跨行业安全协作机制形成",
    ]),
]

for i, (title, color, items) in enumerate(areas):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.3)
    top = Inches(1.5 + row * 2.9)
    card = add_shape(slide, left, top, Inches(5.8), Inches(2.6), BG_MEDIUM, color)
    tb = add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.2), Inches(2.4))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, title, size=20, color=color, bold=True)
    for item in items:
        add_bullet(tb.text_frame, "▸ " + item, size=14, color=LIGHT_GRAY)

# ========== SLIDE 9: Timeline & Outlook ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "时间线与展望", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.05), ACCENT_BLUE)

# Timeline - horizontal
phases = [
    ("2026.3月", "模型泄露", "Mythos存在被数据\n泄露意外曝光", ACCENT_RED),
    ("2026.4.7", "Glasswing启动", "正式发布，12家创始\n伙伴获得访问权限", ACCENT_BLUE),
    ("2026.4.12", "延迟公开发布", "Anthropic宣布推迟\nMythos公开发布", ACCENT_ORANGE),
    ("90天内", "漏洞修复报告", "公开报告已修复漏洞\n和安全改进成果", ACCENT_GREEN),
    ("未来", "Claude Opus演进", "增强安全防护后融入\nClaude Opus主线", WHITE),
]

# Timeline line
add_shape(slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.04), LIGHT_GRAY)

for i, (date, title, desc, color) in enumerate(phases):
    x = Inches(0.8 + i * 2.5)
    # Dot
    dot = add_shape(slide, x + Inches(0.9), Inches(2.65), Inches(0.3), Inches(0.3), color)
    # Date above
    tb = add_text_box(slide, x, Inches(2.0), Inches(2.3), Inches(0.5))
    set_text(tb.text_frame, date, size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Title below
    tb = add_text_box(slide, x, Inches(3.2), Inches(2.3), Inches(0.5))
    set_text(tb.text_frame, title, size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Desc
    tb = add_text_box(slide, x, Inches(3.7), Inches(2.3), Inches(1.2))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, desc, size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key insights
card = add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(2.0), BG_MEDIUM, ACCENT_BLUE)
tb = add_text_box(slide, Inches(1.1), Inches(5.4), Inches(11), Inches(1.7))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "关键洞察", size=20, color=ACCENT_BLUE, bold=True)
insights = [
    "Mythos标志着AI安全能力从量变到质变的转折点 — Opus 4.6在漏洞利用上接近0%成功率，Mythos实现了质的飞跃",
    "\"负责任的有限发布\"模式可能成为未来前沿AI模型的标准做法",
    "定价$25/$125 per M tokens表明Anthropic将此视为高价值商业产品，而非纯粹的安全公益",
    "攻防不对称性在AI时代进一步加剧: 防御者需要修复所有漏洞，攻击者只需找到一个",
]
for insight in insights:
    add_bullet(tb.text_frame, "▸ " + insight, size=13, color=LIGHT_GRAY)

# ========== SLIDE 10: EU AI Act Overview ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

ACCENT_PURPLE = RGBColor(0xAA, 0x55, 0xFF)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "欧盟AI法案 (EU AI Act) 与Mythos/Glasswing", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.05), ACCENT_PURPLE)

# Left: EU AI Act overview
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), BG_MEDIUM, ACCENT_PURPLE)
tb = add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(5.2))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "EU AI Act 关键条款与时间线", size=20, color=ACCENT_PURPLE, bold=True)
add_paragraph(tb.text_frame, "", size=6)

eu_items = [
    ("2024.8.1", "法案正式生效"),
    ("2025.2.2", "禁止性AI实践 & AI素养义务生效"),
    ("2025.8.2", "GPAI模型治理规则和义务生效"),
    ("2026.8.2", "全面适用 — 高风险AI系统规则 & 执行权力启动"),
]
for date, desc in eu_items:
    add_paragraph(tb.text_frame, date, size=16, color=ACCENT_ORANGE, bold=True, space_before=Pt(12))
    add_paragraph(tb.text_frame, desc, size=14, color=LIGHT_GRAY, space_before=Pt(2))

add_paragraph(tb.text_frame, "", size=10)
add_paragraph(tb.text_frame, "Art.55: 系统性风险GPAI义务", size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=4)
art55 = [
    "按标准化协议进行模型评估",
    "开展对抗性测试以识别系统性风险",
    "评估并缓解欧盟层面的系统性风险",
    "追踪和报告严重事件",
    "确保充足的网络安全保护水平",
]
for a in art55:
    add_bullet(tb.text_frame, "▸ " + a, size=13, color=LIGHT_GRAY)

# Right: Mythos classification
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.5), BG_MEDIUM, ACCENT_RED)
tb = add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(5.2))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "Mythos在EU AI Act下的分类", size=20, color=ACCENT_RED, bold=True)
add_paragraph(tb.text_frame, "", size=6)

add_paragraph(tb.text_frame, "通用AI模型 (GPAI)", size=16, color=ACCENT_ORANGE, bold=True, space_before=Pt(10))
add_paragraph(tb.text_frame, "Mythos作为通用语言模型，直接落入\nEU AI Act对GPAI的监管范围", size=13, color=LIGHT_GRAY, space_before=Pt(4))

add_paragraph(tb.text_frame, "具有系统性风险的GPAI", size=16, color=ACCENT_RED, bold=True, space_before=Pt(14))
add_paragraph(tb.text_frame, "训练算力超过10^25 FLOPs即被推定为\n具有系统性风险 — Mythos几乎必然达标", size=13, color=LIGHT_GRAY, space_before=Pt(4))

add_paragraph(tb.text_frame, "双重用途风险", size=16, color=ACCENT_ORANGE, bold=True, space_before=Pt(14))
add_paragraph(tb.text_frame, "自主发现和利用零日漏洞的能力使其\n成为典型的\"双重用途\"AI技术", size=13, color=LIGHT_GRAY, space_before=Pt(4))

add_paragraph(tb.text_frame, "Anthropic的合规立场", size=16, color=ACCENT_GREEN, bold=True, space_before=Pt(14))
add_paragraph(tb.text_frame, "已签署EU GPAI实践准则(Code of Practice)\n加入26家主要AI厂商的合规阵营\n欧盟委员会公开支持Mythos分阶段发布策略", size=13, color=LIGHT_GRAY, space_before=Pt(4))

# ========== SLIDE 11: EU AI Act Impact on Mythos/Glasswing ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "Mythos/Glasswing 对EU AI Act的六大影响", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.05), ACCENT_PURPLE)

impacts = [
    ("1. 加速\"系统性风险\"定义细化", ACCENT_RED,
     "Mythos的自主网络攻击能力远超立法时的预期，\n迫使欧盟重新审视Art.51中\"高影响能力\"\n的具体标准和评估方法论"),
    ("2. 推动网络安全专项条款", ACCENT_BLUE,
     "现有Art.55仅泛泛要求\"充足的网络安全保护\"，\nMythos的出现可能催生针对AI辅助漏洞发现\n和利用的专项监管条款"),
    ("3. 重塑对抗性测试标准", ACCENT_ORANGE,
     "Art.55(1)(a)要求\"对抗性测试\"但缺乏细节，\nMythos级别的安全能力要求全新的红队测试\n协议和评估基准"),
    ("4. 验证\"分阶段发布\"模式", ACCENT_GREEN,
     "Glasswing的限制访问模式与EU AI Act\n\"比例性\"原则高度契合，可能成为\n前沿模型合规发布的参考范式"),
    ("5. 挑战跨境监管框架", ACCENT_PURPLE,
     "Glasswing合作伙伴跨越美欧，漏洞信息\n跨境流动如何符合EU数据保护和\n安全信息共享规则成为新课题"),
    ("6. 倒逼\"双重用途\"AI出口管制", ACCENT_RED,
     "Mythos的军事级网络能力可能触发瓦森纳\n安排下的出口管制讨论，影响AI模型\n的全球分发策略"),
]

for i, (title, color, desc) in enumerate(impacts):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.3)
    top = Inches(1.5 + row * 1.95)
    card = add_shape(slide, left, top, Inches(5.8), Inches(1.75), BG_MEDIUM, color)
    tb = add_text_box(slide, left + Inches(0.25), top + Inches(0.1), Inches(5.3), Inches(1.6))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, title, size=16, color=color, bold=True)
    add_paragraph(tb.text_frame, desc, size=12, color=LIGHT_GRAY, space_before=Pt(4))

# ========== SLIDE 12: EU Regulatory Response & Global Ripple ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
set_text(tb.text_frame, "欧盟监管响应与全球连锁效应", size=32, color=WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(4), Inches(0.05), ACCENT_PURPLE)

# EU Response
card = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0), BG_MEDIUM, ACCENT_PURPLE)
tb = add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(2.7))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "欧盟委员会的响应", size=20, color=ACCENT_PURPLE, bold=True)
add_paragraph(tb.text_frame, "", size=4)
eu_responses = [
    "公开支持Anthropic延迟Mythos公开发布的决定",
    "引用\"大规模网络风险\"作为支持分阶段发布的理由",
    "将Glasswing模式视为AI Act合规的\"良好实践\"",
    "正在评估是否需要为网络安全AI制定补充指南",
    "强调前沿AI的预发布安全评估应成为强制要求",
]
for r in eu_responses:
    add_bullet(tb.text_frame, "▸ " + r, size=13, color=LIGHT_GRAY)

# Compliance challenges
card = add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(3.0), BG_MEDIUM, ACCENT_ORANGE)
tb = add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(2.7))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "合规挑战与矛盾", size=20, color=ACCENT_ORANGE, bold=True)
add_paragraph(tb.text_frame, "", size=4)
challenges = [
    "透明度 vs 安全性: Art.53要求公开文档，但漏洞信息\n  公开可能危及安全",
    "开源义务 vs 限制访问: GPAI Code of Practice鼓励\n  开放，Glasswing模式则高度封闭",
    "跨境数据流: 漏洞信息是否属于需要保护的\"敏感数据\"",
    "责任归属: AI发现的漏洞被利用后，谁承担法律责任",
    "评估成本: 对抗性测试要求推高合规门槛",
]
for c in challenges:
    add_bullet(tb.text_frame, "▸ " + c, size=12, color=LIGHT_GRAY)

# Global ripple effects
card = add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.4), BG_MEDIUM, ACCENT_BLUE)
tb = add_text_box(slide, Inches(1.1), Inches(5.0), Inches(11), Inches(2.1))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "全球连锁效应: EU AI Act的\"布鲁塞尔效应\"", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=4)

global_items = [
    ("美国", "CISA和CAISI已与Anthropic深度对话；Mythos事件加速美国AI安全立法讨论"),
    ("英国", "EU AI Act高风险系统要求将成为英国AI安全研究所(AISI)参考模板"),
    ("日本/韩国", "两国正在制定的AI监管框架将直接借鉴EU AI Act对GPAI系统性风险的条款"),
    ("全球影响", "Mythos/Glasswing验证了\"能力越强、监管越严\"的分级原则，推动全球AI治理共识形成"),
]
for region, desc in global_items:
    add_paragraph(tb.text_frame, region, size=15, color=ACCENT_GREEN, bold=True, space_before=Pt(8))
    add_paragraph(tb.text_frame, desc, size=12, color=LIGHT_GRAY, space_before=Pt(2))

# ========== SLIDE 13: Conclusion ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_shape(slide, Inches(1), Inches(1.5), Inches(0.08), Inches(4.5), ACCENT_BLUE)

tb = add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1))
set_text(tb.text_frame, "总结与建议", size=36, color=WHITE, bold=True)

tb = add_text_box(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(4.5))
tb.text_frame.word_wrap = True
set_text(tb.text_frame, "Claude Mythos + Project Glasswing 的战略意义", size=22, color=ACCENT_BLUE, bold=True)
add_paragraph(tb.text_frame, "", size=8)

conclusions = [
    ("技术层面", "Mythos代表了AI在网络安全领域的里程碑式突破，从\"辅助工具\"进化为\"自主研究员\""),
    ("商业层面", "Anthropic通过\"限制发布+合作伙伴生态\"模式，创造了高壁垒的差异化竞争优势"),
    ("产业层面", "验证了AI在垂直行业创造巨大价值的可能性，符合\"倒金字塔\"健康产业结构的方向"),
    ("安全层面", "攻防博弈进入AI对AI的新阶段，所有组织需重新评估安全策略"),
    ("监管层面", "前沿AI模型的\"负责任发布\"正在形成新的行业范式，政府深度介入成为常态"),
]
for title, desc in conclusions:
    add_paragraph(tb.text_frame, title, size=18, color=ACCENT_ORANGE, bold=True, space_before=Pt(14))
    add_paragraph(tb.text_frame, desc, size=15, color=LIGHT_GRAY, space_before=Pt(2))

# Save
output_path = "/Users/apple/Future_Thoughts/Mythos和Glasswing分析报告.pptx"
prs.save(output_path)
print(f"Report saved to: {output_path}")
