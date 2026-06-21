# -*- coding: utf-8 -*-
"""生成单页PPT：Anthropic云策略"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 配色 ----------
INK   = RGBColor(0x1A, 0x1A, 0x1A)
SUB   = RGBColor(0x5F, 0x5F, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLAY  = RGBColor(0xD9, 0x77, 0x57)   # Anthropic 品牌橙
BG    = RGBColor(0xF7, 0xF4, 0xEF)
AWS   = RGBColor(0xFF, 0x9F, 0x0A)
GCP   = RGBColor(0x42, 0x85, 0xF4)
MSFT  = RGBColor(0x2E, 0x8B, 0x57)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xE3, 0xDD, 0xD3)
FONT  = "PingFang SC"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def set_ea(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def rect(x, y, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def text(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        for (t, sz, col, bold) in ln:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            set_ea(r)
    return tb

# ---------- 背景 ----------
rect(0, 0, prs.slide_width, prs.slide_height, BG)
rect(0, 0, prs.slide_width, Inches(0.18), CLAY)

# ---------- 标题 ----------
text(Inches(0.55), Inches(0.32), Inches(9.8), Inches(0.9),
     [[("Anthropic 云厂商合作策略", 30, INK, True)],
      [("“一模型 · 多云 · 多芯片” —— 以算力采购承诺换取战略投资与产能锁定", 14, SUB, False)]],
     sp_after=4)

tag = rect(Inches(10.55), Inches(0.42), Inches(2.25), Inches(0.55), CLAY)
tf = tag.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "估值约 $3500 亿"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
set_ea(r)

# ---------- 核心主张 ----------
text(Inches(0.55), Inches(1.42), Inches(12.2), Inches(0.5),
     [[("核心打法：", 14, CLAY, True),
       ("Claude 是唯一同时上架全球三大云的前沿模型 —— 最大化分发触达，并在 Trainium / TPU / GPU 间灵活调度保障产能", 14, INK, False)]])

# ---------- 三张卡片 ----------
cards = [
    (AWS, "亚马逊 AWS", "核心云 & 首要训练伙伴",
     [("承诺 10 年投入 $1000 亿+ 采购 AWS 算力", True),
      ("亚马逊累计投资约 $330 亿（最新追加 $250 亿）", False),
      ("Project Rainier：~50 万颗 Trainium2，目标 5GW", False),
      ("分发：Amazon Bedrock · 10 万+ 客户", False)]),
    (GCP, "谷歌 Google Cloud", "TPU 算力主力供给",
     [("2025.10 签约：最多 100 万颗 TPU，2026 超 1GW", True),
      ("谷歌投资 $100 亿 + 按用量追加至 $400 亿", False),
      ("联合 Broadcom：多 GW 新一代算力，2027 起", False),
      ("分发：Vertex AI", False)]),
    (MSFT, "微软 Azure + 英伟达", "补齐 GPU 与第三朵云",
     [("承诺采购 $300 亿 Azure 算力（最高 1GW）", True),
      ("微软投资 $50 亿、英伟达投资 $100 亿", False),
      ("采购英伟达 Grace Blackwell / Vera Rubin", False),
      ("分发：Azure AI Foundry", False)]),
]

cx, cw, gap, cy, ch = Inches(0.55), Inches(3.97), Inches(0.18), Inches(2.05), Inches(3.55)
for i, (color, name, role, bullets) in enumerate(cards):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    rect(x, cy, cw, ch, CARD, line=LINE, line_w=Pt(1))
    rect(x, cy, cw, Inches(0.82), color)
    text(Emu(int(x)+Inches(0.22)), Emu(int(cy)+Inches(0.10)), Emu(int(cw)-Inches(0.4)), Inches(0.7),
         [[(name, 16, WHITE, True)], [(role, 11, WHITE, False)]], sp_after=1)
    lines = []
    for (t, hl) in bullets:
        lines.append([("▪ ", 12, color, True), (t, 11.5, INK if hl else SUB, hl)])
    text(Emu(int(x)+Inches(0.22)), Emu(int(cy)+Inches(0.98)), Emu(int(cw)-Inches(0.42)), Inches(2.4),
         lines, sp_after=7)

# ---------- 底部战略逻辑 ----------
fy = Inches(5.85)
rect(Inches(0.55), fy, Inches(12.23), Inches(1.25), RGBColor(0x2B, 0x2B, 0x2B))
insights = [
    ("产能即护城河", "GW 级长约锁定稀缺算力，对冲供给紧张"),
    ("多芯片去绑定", "Trainium/TPU/GPU 并行，降单一硬件依赖、优化成本"),
    ("分发最大化", "唯一覆盖三云的前沿模型，直达各云存量客户"),
    ("算力换投资", "采购承诺换取巨额投资，形成资本-算力循环"),
]
iw = Inches(3.0)
for i, (h, d) in enumerate(insights):
    x = Emu(int(Inches(0.72)) + i * int(iw))
    text(x, Emu(int(fy)+Inches(0.16)), Emu(int(iw)-Inches(0.18)), Inches(1.0),
         [[(h, 13.5, CLAY, True)], [(d, 10.5, RGBColor(0xDD,0xDD,0xDD), False)]], sp_after=3)

text(Inches(0.55), Inches(7.18), Inches(12.2), Inches(0.3),
     [[("资料来源：Anthropic / Amazon / Google / Microsoft 官方公告及 CNBC、Datacenter Dynamics 报道（2025.10–2026.04）", 8, SUB, False)]])

out = "/Users/apple/Future_Thoughts/Anthropic云策略.pptx"
prs.save(out)
print("saved:", out)
